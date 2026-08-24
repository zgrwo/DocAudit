"""Auto-Fix 引擎 — 自动修复简单格式问题

Inspired by intern's `intern fix` command.
当前支持的修复:
- 字体标准化 (非标准字体 → 指定字体)
- 预留: 对齐修正、文本框位置微调
"""

import logging
import os
import re
import shutil
from pathlib import Path

logger = logging.getLogger(__name__)

# 标题末尾标点去除正则 (模块级预编译，避免每次 fix 调用重复编译)
_TRAILING_PUNCT_RE = re.compile(r"[。，、.!,;；：…—]+$")

# 文本开头的简单符号类项目标记 (模块级预编译)
_BULLET_RE = re.compile(r"^\s*[-*]\s")


class AutoFixer:
    """自动修复文档中的格式问题。

    Usage:
        fixer = AutoFixer(allowed_fonts=["Arial", "微软雅黑"])
        fixer.fix_pptx("input.pptx", "output.pptx")
    """

    def __init__(self, allowed_fonts: list[str] | None = None):
        # 默认对齐 rules.md FMT-001: [微软雅黑, Arial, Noto Sans SC, Calibri]
        self.allowed_fonts = allowed_fonts or ["微软雅黑", "Arial", "Noto Sans SC", "Calibri"]
        self._fix_count = 0

    @property
    def fix_count(self) -> int:
        return self._fix_count

    def fix_pptx(self, source: str | Path, target: str | Path) -> Path:
        """修复 PPTX 文件中的字体问题，输出到 target。

        当前修复:
        1. 将不在 allowed_fonts 中的 Run 字体替换为第一个 allowed font
        2. 修正过小的字号 (< 12pt → 12pt)

        使用原子写入: 先保存到临时文件，成功后再替换目标文件，
        避免保存失败时损坏目标文件。
        """
        source = Path(source)
        target = Path(target)

        from pptx import Presentation
        from pptx.oxml.ns import qn
        from pptx.util import Pt

        prs = Presentation(str(source))

        self._fix_count = 0
        default_font = self.allowed_fonts[0]

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # 修复1: 非标准字体 → 默认字体
                        latin_replaced = False
                        if run.font.name and run.font.name not in self.allowed_fonts:
                            run.font.name = default_font
                            self._fix_count += 1
                            latin_replaced = True

                        # 修复1b: eastAsia (a:ea) 中文字体同步替换
                        # (python-pptx font.name 只写 a:latin，中文显示字体需写 a:ea)
                        try:
                            rPr = run.font._rPr
                            if rPr is not None:
                                ea = rPr.find(qn("a:ea"))
                                if ea is not None:
                                    ea_typeface = ea.get("typeface")
                                    if ea_typeface and ea_typeface not in self.allowed_fonts:
                                        ea.set("typeface", default_font)
                                        self._fix_count += 1
                                elif latin_replaced:
                                    # L1: a:ea 缺失且 latin 被替换 → 创建并写入默认字体
                                    # (python-pptx 1.0.2 未建模 a:ea, 用 lxml 创建, 插在 a:latin 之后)
                                    ea = rPr.makeelement(qn("a:ea"), {})
                                    ea.set("typeface", default_font)
                                    latin = rPr.find(qn("a:latin"))
                                    if latin is not None:
                                        latin.addnext(ea)
                                    else:
                                        rPr.append(ea)
                                    self._fix_count += 1
                        except Exception:
                            pass  # bare-handler-ok — PPTX eastAsia 修复降级，不影响 latin/字号修复

                        # 修复2: 字号过小 → 12pt
                        if run.font.size and run.font.size < 12 * 12700:  # EMU
                            run.font.size = Pt(12)
                            self._fix_count += 1

        # 原子写入: 先写临时文件，成功后再替换
        target = Path(target)
        import tempfile

        tmp_fd, tmp_path = tempfile.mkstemp(
            suffix=target.suffix, prefix="autofix_", dir=target.parent
        )
        os.close(tmp_fd)  # 关闭 fd，避免 Windows 文件锁定
        try:
            prs.save(tmp_path)
            os.replace(tmp_path, str(target))
        except Exception:
            # 清理临时文件
            Path(tmp_path).unlink(missing_ok=True)
            raise

        logger.info("Auto-fix: %d 处字体修复 → %s", self._fix_count, target.name)
        return target

    def fix_docx(self, source: str | Path, target: str | Path) -> Path:
        """修复 DOCX 文件中的字体问题 (原子写入)"""
        source = Path(source)
        target = Path(target)

        from docx import Document
        from docx.oxml.ns import qn
        from docx.shared import Pt

        doc = Document(str(source))

        self._fix_count = 0
        default_font = self.allowed_fonts[0]

        for para in doc.paragraphs:
            for run in para.runs:
                latin_replaced = False
                if run.font.name and run.font.name not in self.allowed_fonts:
                    run.font.name = default_font
                    self._fix_count += 1
                    latin_replaced = True
                # eastAsia 中文字体同步替换 (python-docx font.name 只写 w:ascii/w:hAnsi)
                try:
                    rPr = run._element.rPr
                    if rPr is not None and rPr.rFonts is not None:
                        ea = rPr.rFonts.get(qn("w:eastAsia"))
                        if ea and ea not in self.allowed_fonts:
                            rPr.rFonts.set(qn("w:eastAsia"), default_font)
                            self._fix_count += 1
                        elif not ea and latin_replaced:
                            # L1: w:eastAsia 缺失且 latin 被替换 → 写入默认字体
                            rPr.rFonts.set(qn("w:eastAsia"), default_font)
                            self._fix_count += 1
                    elif latin_replaced:
                        # L1: rFonts 缺失且 latin 被替换 → 创建并写入默认字体
                        run._element.get_or_add_rPr().get_or_add_rFonts().set(
                            qn("w:eastAsia"), default_font
                        )
                        self._fix_count += 1
                except Exception:
                    pass  # bare-handler-ok — DOCX eastAsia 修复降级，不影响 latin/字号修复
                if run.font.size and run.font.size < Pt(12):
                    run.font.size = Pt(12)
                    self._fix_count += 1

        # 原子写入: 先写临时文件，成功后再替换
        import tempfile

        tmp_fd, tmp_path = tempfile.mkstemp(
            suffix=target.suffix, prefix="autofix_", dir=target.parent
        )
        os.close(tmp_fd)  # 关闭 fd，避免 Windows 文件锁定
        try:
            doc.save(tmp_path)
            os.replace(tmp_path, str(target))
        except Exception:
            Path(tmp_path).unlink(missing_ok=True)
            raise

        logger.info("Auto-fix: %d 处修复 → %s", self._fix_count, target.name)
        return target

    def fix_spacing(self, source: str | Path, target: str | Path) -> Path:
        """修复中英文之间的空格 (CJK-Latin spacing)。

        在 CJK 字符与拉丁字符/数字之间自动插入空格。
        支持 PPTX 和 DOCX。使用原子写入：先写临时文件，成功后再替换。
        """
        import tempfile

        from src.text_utils import CJK_LATIN_BOUNDARY, LATIN_CJK_BOUNDARY

        source = Path(source)
        ext = source.suffix.lower()
        target = Path(target)

        # 仅支持 OOXML 格式 (.pptx/.docx)；旧版二进制 .ppt/.doc 不被 python-pptx/docx 支持
        if ext not in (".pptx", ".docx"):
            raise ValueError(f"fix_spacing 仅支持 PPTX/DOCX 格式，不支持: {ext}")

        self._fix_count = 0
        tmp_fd, tmp_path = tempfile.mkstemp(
            suffix=target.suffix, prefix="autofix_spacing_", dir=target.parent
        )
        os.close(tmp_fd)  # 关闭 fd，避免 Windows 文件锁定

        try:
            if ext == ".pptx":
                from pptx import Presentation

                shutil.copy2(source, tmp_path)
                prs = Presentation(str(tmp_path))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                old = run.text
                                new = CJK_LATIN_BOUNDARY.sub(r"\1 \2", old)
                                new = LATIN_CJK_BOUNDARY.sub(r"\1 \2", new)
                                if old != new:
                                    run.text = new
                                    self._fix_count += 1
                prs.save(str(tmp_path))

            elif ext == ".docx":
                from docx import Document

                shutil.copy2(source, tmp_path)
                doc = Document(str(tmp_path))
                for para in doc.paragraphs:
                    for run in para.runs:
                        old = run.text
                        new = CJK_LATIN_BOUNDARY.sub(r"\1 \2", old)
                        new = LATIN_CJK_BOUNDARY.sub(r"\1 \2", new)
                        if old != new:
                            run.text = new
                            self._fix_count += 1
                doc.save(str(tmp_path))

            os.replace(tmp_path, str(target))
        except Exception:
            Path(tmp_path).unlink(missing_ok=True)
            raise

        logger.info("Auto-fix spacing: %d 处空格修复 → %s", self._fix_count, target.name)
        return target

    def fix_element_overflow(self, source: str | Path, target: str | Path) -> Path:
        """修复元素溢出 — 将超出幻灯片边界的元素移回边界内 (inspired by intern fix)。

        PPTX only。使用原子写入。
        """
        import tempfile

        from pptx import Presentation

        source = Path(source)
        target = Path(target)

        if source.suffix.lower() != ".pptx":
            raise ValueError(f"fix_element_overflow 仅支持 PPTX 格式，不支持: {source.suffix}")

        self._fix_count = 0
        tmp_fd, tmp_path = tempfile.mkstemp(
            suffix=".pptx", prefix="autofix_overflow_", dir=target.parent
        )
        os.close(tmp_fd)

        try:
            shutil.copy2(source, tmp_path)
            prs = Presentation(str(tmp_path))
            SW = prs.slide_width  # EMU
            SH = prs.slide_height  # EMU

            for slide in prs.slides:
                for shape in slide.shapes:
                    changed = False
                    # 左溢出 → 移到左边界
                    if shape.left < 0:
                        shape.left = 0
                        changed = True
                    # 右溢出 → 移到右边界内
                    if shape.left + shape.width > SW:
                        shape.left = max(0, SW - shape.width)
                        changed = True
                    # 上溢出 → 移到上边界
                    if shape.top < 0:
                        shape.top = 0
                        changed = True
                    # 下溢出 → 移到底边界内
                    if shape.top + shape.height > SH:
                        shape.top = max(0, SH - shape.height)
                        changed = True

                    if changed:
                        self._fix_count += 1

            prs.save(str(tmp_path))
            os.replace(tmp_path, str(target))
        except Exception:
            Path(tmp_path).unlink(missing_ok=True)
            raise

        logger.info("Auto-fix overflow: %d 处元素位置修复 → %s", self._fix_count, target.name)
        return target

    def fix_title_punctuation(self, source: str | Path, target: str | Path) -> Path:
        """修复标题末尾标点 — 去除标题文本末尾的多余标点符号 (inspired by intern fix)。

        PPTX only。使用原子写入。
        """
        import tempfile

        from pptx import Presentation

        source = Path(source)
        target = Path(target)

        if source.suffix.lower() != ".pptx":
            raise ValueError(f"fix_title_punctuation 仅支持 PPTX 格式，不支持: {source.suffix}")

        self._fix_count = 0
        tmp_fd, tmp_path = tempfile.mkstemp(
            suffix=".pptx", prefix="autofix_title_", dir=target.parent
        )
        os.close(tmp_fd)

        try:
            shutil.copy2(source, tmp_path)
            prs = Presentation(str(tmp_path))

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    # 检测是否为标题形状
                    is_title_shape = False
                    if shape.is_placeholder:
                        try:
                            ph_type = shape.placeholder_format.type
                            # TITLE=1, CENTER_TITLE=3
                            if ph_type in (1, 3):
                                is_title_shape = True
                        except Exception:
                            pass  # bare-handler-ok — 占位符类型读取降级，失败时按非标题形状处理
                    if not is_title_shape and shape.name:
                        if "title" in str(shape.name).lower():
                            is_title_shape = True
                    if not is_title_shape:
                        continue

                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            old = run.text
                            new = _TRAILING_PUNCT_RE.sub("", old)
                            if old != new:
                                run.text = new
                                self._fix_count += 1

            prs.save(str(tmp_path))
            os.replace(tmp_path, str(target))
        except Exception:
            Path(tmp_path).unlink(missing_ok=True)
            raise

        logger.info("Auto-fix title punctuation: %d 处标点修复 → %s", self._fix_count, target.name)
        return target

    def fix_bullet_style(
        self, source: str | Path, target: str | Path, preferred: str = "•"
    ) -> Path:
        """修复项目符号样式 — 将混合项目符号统一为指定样式 (inspired by intern fix)。

        将文本开头的 '-' 或 '*' 替换为 preferred。
        PPTX only。使用原子写入。
        """
        import tempfile

        from pptx import Presentation

        source = Path(source)
        target = Path(target)

        if source.suffix.lower() != ".pptx":
            raise ValueError(f"fix_bullet_style 仅支持 PPTX 格式，不支持: {source.suffix}")

        self._fix_count = 0
        tmp_fd, tmp_path = tempfile.mkstemp(
            suffix=".pptx", prefix="autofix_bullet_", dir=target.parent
        )
        os.close(tmp_fd)

        try:
            shutil.copy2(source, tmp_path)
            prs = Presentation(str(tmp_path))

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            old = run.text
                            new = _BULLET_RE.sub(f"{preferred} ", old)
                            if old != new:
                                run.text = new
                                self._fix_count += 1

            prs.save(str(tmp_path))
            os.replace(tmp_path, str(target))
        except Exception:
            Path(tmp_path).unlink(missing_ok=True)
            raise

        logger.info("Auto-fix bullet style: %d 处项目符号修复 → %s", self._fix_count, target.name)
        return target
