"""PPTX 转换器 — 使用 python-pptx 解析，保留完整格式元数据"""

import logging
from collections import defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Length

from src.converters.base import BaseConverter
from src.models.document import (
    Document,
    DocumentMetadata,
    Page,
    PageElement,
    Paragraph,
    Run,
    TableCell,
)

logger = logging.getLogger(__name__)

# PPTX 内部命名空间
NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


class PptxConverter(BaseConverter):
    """将 PPTX 文件转换为统一 Document 模型。

    保留信息：
    - 每个 Slide 的版式名称 (layout_name)
    - 每个 Shape 的精确位置 (left, top, width, height)
    - 每个 Run 的字体/字号/颜色/加粗/斜体
    - 段落对齐、缩进级别
    - 演讲者备注
    - 表格结构
    - 占位符类型 (title / body / ...)
    """

    def can_handle(self, source_path: str | Path) -> bool:
        ext = Path(source_path).suffix.lower().lstrip(".")
        return ext in ("pptx", "ppt")

    def convert(self, source_path: str | Path) -> Document:
        source_path = Path(source_path)
        logger.info("解析 PPTX: %s", source_path.name)

        try:
            prs = Presentation(str(source_path))
        except Exception as e:
            logger.error("无法打开 PPTX 文件: %s", e)
            raise ValueError(f"无法解析 PPTX 文件: {source_path}") from e

        # --- 元数据 ---
        core_props = prs.core_properties
        metadata = DocumentMetadata(
            title=core_props.title,
            author=core_props.author,
            created=str(core_props.created) if core_props.created else None,
            modified=str(core_props.modified) if core_props.modified else None,
            slide_count=len(prs.slides),
            custom_properties={
                "slide_width": _emu_to_pt(prs.slide_width),
                "slide_height": _emu_to_pt(prs.slide_height),
            },
        )

        # --- 遍历 Slide ---
        pages: list[Page] = []

        for slide_idx, slide in enumerate(prs.slides):
            # 版式名称
            layout_name = None
            try:
                layout_name = slide.slide_layout.name
            except Exception:
                pass  # bare-handler-ok — 版式名读取降级，失败时保留 None

            # 演讲者备注
            notes_text = None
            try:
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass  # bare-handler-ok — 备注提取降级，失败时保留 None

            elements: list[PageElement] = []

            # 按 z-order 遍历 shape
            for shape in slide.shapes:
                element = self._convert_shape(shape)
                if element is not None:
                    elements.append(element)

            pages.append(
                Page(
                    index=slide_idx,
                    elements=elements,
                    layout_name=layout_name,
                    slide_number=slide_idx + 1,
                    notes=notes_text or None,
                )
            )

        logger.info(
            "PPTX 解析完毕: %d slides, %d 总元素", len(pages), sum(len(p.elements) for p in pages)
        )

        return Document(
            source_path=str(source_path),
            format="pptx",
            metadata=metadata,
            pages=pages,
        )

    # ── Shape 转换 ───────────────────────────────────────────

    def _convert_shape(self, shape: Any) -> PageElement | None:
        """将单个 Shape 转换为 PageElement。
        单个 shape 转换失败不影响同一 Slide 的其他 shape。
        """
        try:
            shape_type = shape.shape_type

            if shape_type == MSO_SHAPE_TYPE.GROUP:
                return self._convert_group(shape)

            if shape.has_table:
                return self._convert_table(shape)

            if shape.has_text_frame:
                return self._convert_text_frame(shape)

            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                return self._convert_image(shape)

            if shape_type == MSO_SHAPE_TYPE.CHART:
                return self._convert_chart(shape)

            return None
        except (ValueError, TypeError, AttributeError, KeyError, IndexError) as e:
            # 安全获取 shape 标识 — 避免 except 块内二次异常 (shape 属性访问可能同样失败)
            try:
                shape_label = f"{shape.name} (type={shape.shape_type})"
            except Exception:
                shape_label = f"<unknown shape: {type(shape).__name__}>"
            logger.warning("Shape 转换失败: %s: %s", shape_label, e)
            return None

    def _convert_text_frame(self, shape: Any) -> PageElement:
        """转换文本框/占位符"""
        tf = shape.text_frame
        paragraphs: list[Paragraph] = []

        for para in tf.paragraphs:
            runs: list[Run] = []
            # 行距: 固定行距返回 Centipoints (Length, int 子类)，转为 pt 以满足 float|None 契约 (F5)
            line_spacing = para.line_spacing if para.line_spacing else None
            if isinstance(line_spacing, Length):
                line_spacing = line_spacing.pt
            for run_elem in para.runs:
                font = run_elem.font
                # 安全获取颜色
                try:
                    font_color = (
                        _rgb_to_hex(font.color.rgb) if font.color and font.color.rgb else None
                    )
                except Exception:
                    font_color = None

                # eastAsia 中文字体: a:rPr 下的 a:ea typeface
                # (python-pptx font.name 只读 a:latin，不含中文显示字体)
                font_name_east_asia = _ea_typeface(run_elem)

                runs.append(
                    Run(
                        text=run_elem.text,
                        font_name=font.name,
                        font_name_east_asia=font_name_east_asia,
                        font_size=font.size / 12700 if font.size else None,  # EMU → pt
                        bold=font.bold,
                        italic=font.italic,
                        underline=font.underline,
                        color=font_color,
                    )
                )

            paragraphs.append(
                Paragraph(
                    text=para.text,
                    runs=runs,
                    level=para.level,
                    alignment=_pptx_alignment(para.alignment),
                    space_before=para.space_before.pt if para.space_before else None,
                    space_after=para.space_after.pt if para.space_after else None,
                    line_spacing=line_spacing,
                )
            )

        # 判断是否为占位符
        is_placeholder = shape.is_placeholder
        ph_type = None
        if is_placeholder:
            try:
                ph_type = shape.placeholder_format.type
            except Exception:
                ph_type = None  # 占位符格式损坏，回退为通用占位符

        return PageElement(
            type="text_frame",
            paragraphs=paragraphs,
            left=_emu_to_pt(shape.left) if shape.left is not None else None,
            top=_emu_to_pt(shape.top) if shape.top is not None else None,
            width=_emu_to_pt(shape.width) if shape.width is not None else None,
            height=_emu_to_pt(shape.height) if shape.height is not None else None,
            shape_name=shape.name,
            is_title=(is_placeholder and ph_type in (1, 3)),  # TITLE=1, CENTER_TITLE=3
            is_body=(
                is_placeholder and ph_type in (2, 4, 6)
            ),  # BODY=2, SUBTITLE=4, VERTICAL_BODY=6
            is_placeholder=is_placeholder,
        )

    def _convert_table(self, shape: Any) -> PageElement | None:
        """转换表格"""
        table = shape.table
        if len(table.rows) == 0:
            return None  # 空表格 — 无有意义内容
        cells: list[TableCell] = []
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                tf = cell.text_frame
                cell_text = tf.text.strip()

                # 取第一个 run 的格式作为单元格格式
                font_name = None
                font_size = None
                font_color = None
                if tf.paragraphs and tf.paragraphs[0].runs:
                    first_run = tf.paragraphs[0].runs[0]
                    font_name = first_run.font.name
                    if first_run.font.size:
                        font_size = first_run.font.size / 12700
                    try:
                        if first_run.font.color and first_run.font.color.rgb:
                            font_color = _rgb_to_hex(first_run.font.color.rgb)
                    except Exception:
                        font_color = None  # bare-handler-ok — 主题色等无法取 RGB，降级

                # 单元格底色 (仅 solid 纯色填充；渐变/图案/无填充 → None)
                fill_color = None
                try:
                    if cell.fill.type == MSO_FILL_TYPE.SOLID:
                        fill_color = _rgb_to_hex(cell.fill.fore_color.rgb)
                except Exception:
                    fill_color = None  # bare-handler-ok — 异常填充结构，降级

                cells.append(
                    TableCell(
                        text=cell_text,
                        row=row_idx,
                        col=col_idx,
                        font_name=font_name,
                        font_size=font_size,
                        fill_color=fill_color,
                        font_color=font_color,
                    )
                )

        # 按行分组 (单次遍历)
        row_map: dict[int, list[TableCell]] = defaultdict(list)
        for c in cells:
            row_map[c.row].append(c)
        rows = [row_map[r] for r in sorted(row_map)]

        return PageElement(
            type="table",
            tables=rows,
            left=_emu_to_pt(shape.left) if shape.left is not None else None,
            top=_emu_to_pt(shape.top) if shape.top is not None else None,
            width=_emu_to_pt(shape.width) if shape.width is not None else None,
            height=_emu_to_pt(shape.height) if shape.height is not None else None,
            shape_name=shape.name,
        )

    def _convert_image(self, shape: Any) -> PageElement:
        """转换图片 — 仅保留格式信息 (image_ext)，不把整图 blob 载入内存 (P1-5)"""
        image_ext = "unknown"
        if shape.image and shape.image.content_type:
            image_ext = shape.image.content_type.split("/")[-1]
            if image_ext.startswith("x-"):
                image_ext = image_ext[2:]  # "x-emf" → "emf"
        return PageElement(
            type="image",
            left=_emu_to_pt(shape.left) if shape.left is not None else None,
            top=_emu_to_pt(shape.top) if shape.top is not None else None,
            width=_emu_to_pt(shape.width) if shape.width is not None else None,
            height=_emu_to_pt(shape.height) if shape.height is not None else None,
            shape_name=shape.name,
            image_ext=image_ext,
        )

    def _convert_chart(self, shape: Any) -> PageElement:
        """转换图表 — 仅保留 chart_type，不读内嵌 Excel 整包 blob (P1-5)"""
        chart_type = None
        try:
            chart = shape.chart
            chart_type = str(chart.chart_type) if chart.chart_type else None
        except Exception:
            pass  # bare-handler-ok — 图表类型提取降级，失败时保留 None，不阻塞审查

        return PageElement(
            type="chart",
            left=_emu_to_pt(shape.left) if shape.left is not None else None,
            top=_emu_to_pt(shape.top) if shape.top is not None else None,
            width=_emu_to_pt(shape.width) if shape.width is not None else None,
            height=_emu_to_pt(shape.height) if shape.height is not None else None,
            shape_name=shape.name,
            chart_type=chart_type,
        )

    def _convert_group(self, shape: Any) -> PageElement:
        """转换群组 — 递归解析子形状"""
        child_elements: list[PageElement] = []
        for child in shape.shapes:
            element = self._convert_shape(child)
            if element is not None:
                child_elements.append(element)

        # 保留子元素 — 审计器可递归检查 children
        return PageElement(
            type="group",
            children=child_elements,
            left=_emu_to_pt(shape.left) if shape.left is not None else None,
            top=_emu_to_pt(shape.top) if shape.top is not None else None,
            width=_emu_to_pt(shape.width) if shape.width is not None else None,
            height=_emu_to_pt(shape.height) if shape.height is not None else None,
            shape_name=shape.name,
        )

    # ── 图表数据提取 ─────────────────────────────────────────
    # (P1-5: 不再读取内嵌 Excel 整包 blob — 模型仅保留 chart_type)


# ── 工具函数 ────────────────────────────────────────────────


def _ea_typeface(run_elem) -> str | None:
    """只读提取 run 的 a:ea typeface (F4)。

    不使用 run.font._rPr (get_or_add 会凭空创建空 a:rPr)，改走只读
    run._r.find(qn("a:rPr"))；a:rPr 或 a:ea 不存在时返回 None，保持现有行为。
    """
    try:
        rPr = run_elem._r.find(qn("a:rPr"))
        if rPr is not None:
            ea = rPr.find(qn("a:ea"))
            if ea is not None:
                return ea.get("typeface")
    except Exception:
        pass  # bare-handler-ok — ea 字体提取降级，失败时保留 None
    return None


def _emu_to_pt(emu_value) -> float | None:
    """EMU (English Metric Unit) → points (1 pt = 12700 EMU)"""
    if emu_value is None:
        return None
    if isinstance(emu_value, (int, float)):
        return round(emu_value / 12700, 2)
    return round(int(emu_value) / 12700, 2)


def _rgb_to_hex(rgb) -> str | None:
    """RGB 颜色 → hex 字符串"""
    if rgb is None:
        return None
    try:
        return str(rgb)
    except Exception:
        return None


def _pptx_alignment(align) -> str | None:
    """python-pptx 对齐枚举 → 字符串"""
    if align is None:
        return None
    mapping = {
        1: "left",
        2: "center",
        3: "right",
        4: "justify",
    }
    return mapping.get(int(align), None)
