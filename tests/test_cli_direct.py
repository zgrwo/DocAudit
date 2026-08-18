"""cli.py 直接单元测试 (M12) — 把 cli.py 从 coverage 0% 拉起 (关键路径)。

覆盖:
- audit_file: 临时文件 → (Document, list[AuditFinding]) 结构；不存在 → FileNotFoundError；
  不支持格式 → ValueError
- print_summary: 输出含 [SUMMARY] 与 ERROR/WARN/INFO 计数
- main --fix 链路: fixture PPTX 直测 → *_fixed.pptx 输出存在且字体修复生效
"""

import sys
from pathlib import Path

import pytest

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from src.cli import audit_file, main, print_summary  # noqa: E402
from src.models.document import Document  # noqa: E402
from src.models.finding import AuditFinding, FindingSeverity, FindingType  # noqa: E402

RULES = str(PROJECT_ROOT / "rules.md")

# 与 rules.md FMT-001 允许字体一致 (改动 rules.md 需同步此清单)
ALLOWED_FONTS = ["微软雅黑", "Arial", "Noto Sans SC", "Calibri"]


def _needs_fix_pptx(path: Path) -> None:
    """生成含非标准字体 (宋体) + 过小字号 (10pt) 的 PPTX (--fix 可修复)。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    run = txBox.text_frame.paragraphs[0].add_run()
    run.text = "FinFET 技术概述"
    run.font.name = "宋体"  # 非标准字体
    run.font.size = Pt(10)  # 过小字号
    prs.save(str(path))


class TestAuditFile:
    """audit_file: 返回 (doc, findings) 结构 + 错误路径。"""

    def test_returns_doc_and_findings(self, tmp_path):
        pptx = tmp_path / "audit.pptx"
        _needs_fix_pptx(pptx)
        doc, findings = audit_file(str(pptx), rules_path=RULES)
        assert isinstance(doc, Document)
        assert isinstance(findings, list)
        assert all(isinstance(f, AuditFinding) for f in findings)
        assert doc.format == "pptx"
        assert doc.source_path == str(pptx)

    def test_missing_file_raises(self, tmp_path):
        with pytest.raises(FileNotFoundError):
            audit_file(str(tmp_path / "nope.pptx"), rules_path=RULES)

    def test_unsupported_format_raises(self, tmp_path):
        exe = tmp_path / "x.exe"
        exe.write_bytes(b"x")
        with pytest.raises(ValueError, match="不支持的文件格式"):
            audit_file(str(exe), rules_path=RULES)


class TestPrintSummary:
    """print_summary: 输出含 SUMMARY 与 ERROR/WARN/INFO 计数。"""

    def test_summary_counts(self, capsys):
        findings = [
            AuditFinding(type=FindingType.FORMAT, severity=FindingSeverity.ERROR, message="e1"),
            AuditFinding(type=FindingType.FORMAT, severity=FindingSeverity.WARNING, message="w1"),
            AuditFinding(type=FindingType.FORMAT, severity=FindingSeverity.INFO, message="i1"),
            AuditFinding(type=FindingType.FORMAT, severity=FindingSeverity.ERROR, message="e2"),
        ]
        print_summary(findings)
        out = capsys.readouterr().out
        assert "[SUMMARY]" in out
        assert "[ERROR] 2" in out
        assert "[WARN]  1" in out
        assert "[INFO]  1" in out


class TestMainFixChain:
    """main() --fix 链路: 修复生效 + 输出文件存在 (fixture PPTX 直测)。"""

    def test_fix_chain_writes_fixed_file(self, tmp_path):
        pptx = tmp_path / "needs_fix.pptx"
        _needs_fix_pptx(pptx)
        old_argv = sys.argv
        try:
            sys.argv = ["docaudit", str(pptx), "--rules", RULES, "--fix"]
            with pytest.raises(SystemExit):
                main()
        finally:
            sys.argv = old_argv

        out = tmp_path / "needs_fix_fixed.pptx"
        assert out.exists(), "修复输出文件应存在"

        from pptx import Presentation

        prs = Presentation(str(out))
        runs = [
            r
            for s in prs.slides
            for sh in s.shapes
            if sh.has_text_frame
            for p in sh.text_frame.paragraphs
            for r in p.runs
        ]
        assert runs, "修复后文件应有文本 run"
        for run in runs:
            if run.font.name:
                assert run.font.name in ALLOWED_FONTS, (
                    f"字体应被替换为允许字体, got {run.font.name}"
                )
            if run.font.size:
                assert run.font.size >= 12 * 12700, f"字号应修正到 >=12pt, got {run.font.size}"
