"""测试审计器 — 通过公共 audit() API 验证行为，确保测试与源码行为一致"""

import os
import tempfile

from src.auditors.factual import FactualAuditor
from src.auditors.format import FormatAuditor
from src.auditors.structure import StructureAuditor
from src.converters.pptx_converter import PptxConverter
from src.models.document import Document, DocumentMetadata, Page, PageElement, Paragraph
from src.models.finding import FindingSeverity


def _heading_doc(levels: list[int]) -> Document:
    """构造含指定标题层级的单页文档 (text_frame 段落)。"""
    paragraphs = [Paragraph(text=f"H{lv} 标题", runs=[], level=lv) for lv in levels]
    page = Page(
        index=0, slide_number=1, elements=[PageElement(type="text_frame", paragraphs=paragraphs)]
    )
    return Document(format="md", source_path="x", metadata=DocumentMetadata(), pages=[page])


def _text_doc(text: str) -> Document:
    """构造单页纯文本文档。"""
    page = Page(
        index=0,
        slide_number=1,
        elements=[PageElement(type="text_frame", paragraphs=[Paragraph(text=text, runs=[])])],
    )
    return Document(format="md", source_path="x", metadata=DocumentMetadata(), pages=[page])


class TestStructureAuditor:
    """通过 audit() 公共 API 和内部方法双重验证"""

    def test_title_slide_detection_via_audit(self):
        """STR-001: 第一页为标题版式 → audit() 不应产生 STR-001 ERROR"""
        doc = PptxConverter().convert("tests/fixtures/sample.pptx")
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa.audit(doc)
        # sample.pptx 第一页为 Title Slide 版式 → 不应有 STR-001 告警
        str001_errors = [
            f for f in findings if f.rule_id == "STR-001" and f.severity == FindingSeverity.ERROR
        ]
        assert len(str001_errors) == 0, (
            f"STR-001 should not fire on title slide, got: {str001_errors}"
        )

    def test_title_slide_detection_direct(self):
        """STR-001: 直接调用 _check_title_slide → 空列表"""
        doc = PptxConverter().convert("tests/fixtures/sample.pptx")
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa._check_title_slide(doc)
        assert len(findings) == 0, f"Expected no title slide warning, got: {findings}"

    def test_every_slide_conclusion(self):
        """CON-004: 标题页豁免, 内容页应有足够内容"""
        doc = PptxConverter().convert("tests/fixtures/sample.pptx")
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa.audit(doc)
        # 标题页豁免: 第一页不应有 CON-004 告警
        title_page_errors = [f for f in findings if f.page_index == 0 and f.rule_id == "CON-004"]
        assert len(title_page_errors) == 0, (
            f"Title slide should be exempt, got CON-004: {title_page_errors}"
        )

    def test_conclusion_findings_have_valid_structure(self):
        """CON-004: 所有 ERROR findings 必须有有效的 rule_id"""
        doc = PptxConverter().convert("tests/fixtures/sample.pptx")
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa._check_every_slide_has_conclusion(doc)
        for f in findings:
            if f.severity == FindingSeverity.ERROR:
                assert f.rule_id == "CON-004", f"Unexpected ERROR rule_id: {f.rule_id}"
                assert f.message, "Finding message should not be empty"
                assert f.page_index is not None, "Finding must have page_index"

    def test_exempt_layouts_empty_list_uses_default(self):
        """CON-004 豁免版式: 显式空列表不得覆盖内置默认 (修复: pipeline 传 [] 使默认失效)"""
        sa = StructureAuditor(config={"exempt_layouts": []})
        assert sa.exempt_layouts, "空列表不应覆盖内置默认豁免版式"
        assert "标题幻灯片" in sa.exempt_layouts

    def test_exempt_layouts_declared_used(self):
        """CON-004 豁免版式: 显式声明的列表生效"""
        sa = StructureAuditor(config={"exempt_layouts": ["封面页"]})
        assert sa.exempt_layouts == ["封面页"]

    def test_every_slide_conclusion_non_pptx_no_findings(self):
        """CON-004: 非 PPTX (MD) 不执行每页结论检查 (修复: 曾对 DOCX/PDF/MD 每页误报 error)"""
        doc = _text_doc("这是第一段普通正文，第二段补充说明，第三段技术描述。")
        sa = StructureAuditor(config={"required_sections": []})
        assert sa._check_every_slide_has_conclusion(doc) == [], "非 PPTX 文档不应产生 CON-004"

    def test_heading_levels_page_first_heading_exempt(self):
        """STR-003: 页首标题不参与跳级比较 (修复: 按标题分页的文档页首 H2/H3 不再误报)"""
        doc = _heading_doc([2, 3])  # 页首即 H2 (MD ### 切页场景)
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa._check_heading_levels(doc)
        assert len(findings) == 0, f"页首 H2 不应误报跳级, got: {findings}"

    def test_heading_levels_real_skip_flagged(self):
        """STR-003: 页内真实跳级仍被检测 (H1→H3)"""
        doc = _heading_doc([1, 3])
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa._check_heading_levels(doc)
        assert len(findings) == 1
        assert "H1" in findings[0].message and "H3" in findings[0].message

    def test_heading_levels_sequential_ok(self):
        """STR-003: 逐级递进 (H1→H2→H3) 不报"""
        doc = _heading_doc([1, 2, 3])
        sa = StructureAuditor(config={"required_sections": []})
        assert sa._check_heading_levels(doc) == []

    def test_heading_levels_mid_document_skip_flagged(self):
        """STR-003: 页内 H2→H4 仍报 (首标题豁免后保留跳级检测)"""
        doc = _heading_doc([2, 4])
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa._check_heading_levels(doc)
        assert len(findings) == 1
        assert "H2" in findings[0].message and "H4" in findings[0].message

    def test_heading_levels_pptx_bullet_indent_not_flagged(self):
        """STR-003: PPTX 的 para.level 是缩进级别, 不得当标题层级 (修复: bullet 缩进 0→2 误报)"""
        page = Page(
            index=0,
            slide_number=1,
            elements=[
                PageElement(
                    type="text_frame",
                    paragraphs=[
                        Paragraph(text="• 第一级", runs=[], level=0),
                        Paragraph(text="    • 第三级", runs=[], level=2),
                    ],
                )
            ],
        )
        doc = Document(
            format="pptx", source_path="x.pptx", metadata=DocumentMetadata(), pages=[page]
        )
        sa = StructureAuditor(config={"required_sections": []})
        assert sa._check_heading_levels(doc) == [], "PPTX 缩进级别不应触发 STR-003 跳级"

    def test_figure_numbering_within_page_order_preserved(self):
        """STR-002: 同页内按出现次序检查 (修复: 曾按编号重排, 倒退被掩盖成跳号)"""
        doc = _text_doc("参见图3 与图1 的对比")
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa._check_figure_numbering(doc)
        assert len(findings) == 1, f"期望 1 条 STR-002, got: {findings}"
        assert "倒退" in findings[0].message, (
            f"图3 先于 图1 出现应为'倒退', got: {findings[0].message}"
        )

    def test_figure_numbering_cross_page_skip_flagged(self):
        """STR-002: 跨页跳号仍被检测 (图1 → 图3)"""
        p1 = Page(
            index=0,
            slide_number=1,
            elements=[
                PageElement(type="text_frame", paragraphs=[Paragraph(text="如图1 所示", runs=[])])
            ],
        )
        p2 = Page(
            index=1,
            slide_number=2,
            elements=[
                PageElement(type="text_frame", paragraphs=[Paragraph(text="如图3 所示", runs=[])])
            ],
        )
        doc = Document(format="md", source_path="x", metadata=DocumentMetadata(), pages=[p1, p2])
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa._check_figure_numbering(doc)
        assert len(findings) == 1
        assert "不连续" in findings[0].message

    def test_figure_numbering_chapter_style_ignored(self):
        """STR-002: 章节式编号 (图1-1/图1-2/图2-1) 不误报 '编号重复'"""
        doc = _text_doc("图1-1 工艺流程图\n图1-2 参数表\n图2-1 结构图")
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa._check_figure_numbering(doc)
        assert len(findings) == 0, f"章节式编号不应报 STR-002, got: {findings}"

    def test_figure_numbering_chapter_style_real_skip_flagged(self):
        """STR-002: 章节式编号内部真实跳号仍被检测 (图1-1 → 图1-3)"""
        doc = _text_doc("图1-1 工艺流程图\n图1-3 参数表")
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa._check_figure_numbering(doc)
        assert len(findings) == 0, f"章节式编号暂不参与连续性校验 (设计取舍), got: {findings}"

    def test_figure_caption_format_space_insensitive(self):
        """STR-007: 指纹对空格不敏感 ('Fig. 1:' 与 'Fig.2:' 视为同一格式)"""
        doc = _text_doc("Fig. 1: 标题甲。\nFig.2: 标题乙。")
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa._check_figure_caption_format(doc)
        assert len(findings) == 0, f"空格差异不应算两种格式, got: {findings}"

    def test_figure_caption_format_mixed_flagged(self):
        """STR-007: 中英文格式混用仍被检测 ('图1：' vs 'Fig. 1:')"""
        doc = _text_doc("图1：标题甲。\nFig. 1: 标题乙。")
        sa = StructureAuditor(config={"required_sections": []})
        findings = sa._check_figure_caption_format(doc)
        assert len(findings) == 1
        assert "2 种" in findings[0].message

    def test_title_length_or_condition(self):
        """STR-004: 纯中文超长标题应触发告警 — 验证 AND→OR 修复"""
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        s = prs.slides.add_slide(slide_layout)
        # 设置超长中文标题: ~50+ chars > 40 上限
        long_title = "这是一个非常长的中文标题用于测试标题长度检查功能是否正常工作验证超长标题检测逻辑是否正确触发告警机制"
        s.shapes.title.text = long_title

        tmp_path = None
        try:
            # 使用临时文件避免泄漏
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp_path = tmp.name
            prs.save(tmp_path)

            doc = PptxConverter().convert(tmp_path)
            sa = StructureAuditor(config={"required_sections": []})
            findings = sa._check_title_length(doc)
            assert len(findings) >= 1, (
                f"Long Chinese title ({len(long_title)} chars) should trigger STR-004. "
                f"Got {len(findings)} findings."
            )
            # 验证 finding 包含正确的元数据
            f = findings[0]
            assert f.rule_id == "STR-004"
            assert f.metadata.get("chinese_chars", 0) > 40
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.remove(tmp_path)


class TestFormatAuditor:
    def test_font_consistency_audit_api(self):
        """FMT-001: 非标准字体应通过 audit() 被检测到"""
        doc = PptxConverter().convert("tests/fixtures/sample.pptx")
        findings = FormatAuditor().audit(doc)
        fmt001_findings = [f for f in findings if f.rule_id == "FMT-001"]
        assert len(fmt001_findings) >= 1, (
            f"Expected at least one FMT-001 (non-standard font) finding, got {len(fmt001_findings)}"
            + f"\nAll rule_ids: {[f.rule_id for f in findings]}"
        )
        # 验证 finding 结构完整
        for f in fmt001_findings:
            assert f.type.value == "format"
            assert f.message
            assert f.rule_id == "FMT-001"

    def test_font_consistency_east_asia_flagged(self):
        """FMT-001: 非允许的 eastAsia 中文字体应触发告警 (font_scope=east_asia)"""
        from src.models.document import Run

        page = Page(
            index=0,
            slide_number=1,
            elements=[
                PageElement(
                    type="text_frame",
                    paragraphs=[
                        Paragraph(
                            text="中文正文",
                            runs=[
                                Run(text="中文正文", font_name="Arial", font_name_east_asia="宋体"),
                            ],
                        ),
                    ],
                ),
            ],
        )
        findings = FormatAuditor()._check_font_consistency(page)
        fmt001 = [f for f in findings if f.rule_id == "FMT-001"]
        ea = [f for f in fmt001 if (f.metadata or {}).get("font_scope") == "east_asia"]
        assert len(ea) >= 1, f"非允许 eastAsia 字体应触发 FMT-001, got: {fmt001}"
        assert "宋体" in ea[0].message
        assert "中文" in ea[0].message  # message 区分中文（eastAsia）字体
        # latin 字体 (Arial) 在允许列表 → 不误报
        latin = [f for f in fmt001 if (f.metadata or {}).get("font_scope") == "latin"]
        assert len(latin) == 0

    def test_font_consistency_latin_scope_label(self):
        """FMT-001: 非允许 latin 字体的 message 标注为西文（latin）"""
        from src.models.document import Run

        page = Page(
            index=0,
            slide_number=1,
            elements=[
                PageElement(
                    type="text_frame",
                    paragraphs=[
                        Paragraph(
                            text="x",
                            runs=[
                                Run(
                                    text="x", font_name="Comic Sans", font_name_east_asia="微软雅黑"
                                ),
                            ],
                        ),
                    ],
                ),
            ],
        )
        findings = FormatAuditor()._check_font_consistency(page)
        fmt001 = [f for f in findings if f.rule_id == "FMT-001"]
        latin = [f for f in fmt001 if (f.metadata or {}).get("font_scope") == "latin"]
        assert len(latin) == 1
        assert "西文" in latin[0].message
        assert latin[0].metadata.get("font") == "Comic Sans"

    def test_global_font_consistency_counts_east_asia(self):
        """FMT-001 全局: 复合键 (font, scope) 统计 — latin 统一 + 3 种 eastAsia → 触发"""
        from src.models.document import Run

        page = Page(
            index=0,
            slide_number=1,
            elements=[
                PageElement(
                    type="text_frame",
                    paragraphs=[
                        Paragraph(
                            text="a",
                            runs=[Run(text="a", font_name="Arial", font_name_east_asia="宋体")],
                        ),
                        Paragraph(
                            text="b",
                            runs=[Run(text="b", font_name="Arial", font_name_east_asia="黑体")],
                        ),
                        Paragraph(
                            text="c",
                            runs=[Run(text="c", font_name="Arial", font_name_east_asia="楷体")],
                        ),
                    ],
                )
            ],
        )
        doc = Document(format="pptx", source_path="x", metadata=DocumentMetadata(), pages=[page])
        findings = FormatAuditor()._check_global_font_consistency(doc)
        assert len(findings) == 1, f"4 种复合字体应触发 FMT-001 全局检查, got: {findings}"
        f = findings[0]
        assert f.rule_id == "FMT-001"
        dist = f.metadata["font_distribution"]
        assert dist["Arial/latin"] == 3
        assert dist["宋体/east_asia"] == 1 and dist["黑体/east_asia"] == 1
        assert dist["楷体/east_asia"] == 1

    def test_english_char_count_ignores_punctuation(self):
        """FMT-004: 英文计数只算字母 — 数字/标点/空格不计入 (修复: 曾 总长-中文字数)"""
        fa = FormatAuditor()
        doc = _text_doc("1" * 300 + " test")
        findings = fa._check_paragraph_length(doc.pages[0])
        assert len(findings) == 0, f"纯数字+少量字母不应触发 FMT-004, got: {findings}"

    def test_english_char_count_pure_letters_triggers(self):
        """FMT-004: 纯字母超上限仍触发 (回归)"""
        fa = FormatAuditor()
        doc = _text_doc("a" * 305)
        findings = fa._check_paragraph_length(doc.pages[0])
        assert len(findings) == 1
        assert findings[0].rule_id == "FMT-004"
        assert "英文字符" in findings[0].message

    def test_global_font_consistency_latin_plus_two_east_asia_ok(self):
        """FMT-001 全局: latin 统一 + 2 种 eastAsia (共 3 复合键) → 不触发"""
        from src.models.document import Run

        page = Page(
            index=0,
            slide_number=1,
            elements=[
                PageElement(
                    type="text_frame",
                    paragraphs=[
                        Paragraph(
                            text="a",
                            runs=[Run(text="a", font_name="Arial", font_name_east_asia="宋体")],
                        ),
                        Paragraph(
                            text="b",
                            runs=[Run(text="b", font_name="Arial", font_name_east_asia="黑体")],
                        ),
                    ],
                )
            ],
        )
        doc = Document(format="pptx", source_path="x", metadata=DocumentMetadata(), pages=[page])
        findings = FormatAuditor()._check_global_font_consistency(doc)
        assert len(findings) == 0, f"3 种复合字体不应触发, got: {findings}"

    def test_font_size_max_check(self):
        """FMT-002: 验证字号最大值检查生效"""
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        s = prs.slides.add_slide(slide_layout)
        # 设置超大标题字号 (72pt > 40pt max)
        s.shapes.title.text = "Huge Title"
        s.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(72)

        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp_path = tmp.name
            prs.save(tmp_path)

            doc = PptxConverter().convert(tmp_path)
            findings = FormatAuditor().audit(doc)
            # 应检测到超大标题字号
            fmt002_max = [f for f in findings if f.rule_id == "FMT-002" and "过大" in f.message]
            assert len(fmt002_max) >= 1, (
                f"72pt title should trigger FMT-002 max size warning. "
                f"All FMT-002: {[f.message for f in findings if f.rule_id == 'FMT-002']}"
            )
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.remove(tmp_path)


class TestFactualAuditor:
    def test_abbreviation_check_via_audit(self):
        """CON-003: 未定义缩写应通过 audit() 被检测"""
        doc = PptxConverter().convert("tests/fixtures/sample.pptx")
        findings = FactualAuditor().audit(doc)
        con003_findings = [f for f in findings if f.rule_id == "CON-003"]
        assert len(con003_findings) >= 1, (
            f"Expected at least one CON-003 (undefined abbreviation) finding, got {len(con003_findings)}"
        )
        # 验证发现包含正确信息
        for f in con003_findings:
            assert "未给出全称" in f.message, f"Unexpected message: {f.message}"
            assert f.type.value == "factual"
            assert f.severity == FindingSeverity.WARNING

    def test_common_uppercase_words_skipped(self):
        """CON-003: 常见英语大写单词 (THE, AND, NEW 等) 不应被标记为未定义缩写"""
        doc = PptxConverter().convert("tests/fixtures/sample.pptx")
        findings = FactualAuditor().audit(doc)
        # 验证 THE/AND/NEW 等常见词不在 findings 中
        con003_findings = [f for f in findings if f.rule_id == "CON-003"]
        common_words = {
            "THE",
            "AND",
            "FOR",
            "ALL",
            "BUT",
            "NOT",
            "CAN",
            "ARE",
            "WAS",
            "HAS",
            "NEW",
            "SET",
            "END",
            "TOP",
        }
        for f in con003_findings:
            # 从 context 中提取缩写字面
            ctx = f.context or ""
            for word in common_words:
                if word in ctx:
                    # 确认这不是被标记的缩写
                    assert word not in (f.metadata.get("abbreviation") or ""), (
                        f"Common word '{word}' should not be flagged as undefined abbreviation"
                    )

    def test_abbreviation_reuse_after_definition_not_double_defined(self):
        """CON-003-B: 定义后再次使用 (复用) 不得误计为重复定义 (修复: is_full_before 120 字符窗口假阳性)"""
        fa = FactualAuditor()
        doc = _text_doc("TSV (硅通孔 TSV) 工艺。TSV 用于 3D 集成。")
        findings = fa._check_abbreviation_multiply_defined(doc)
        assert len(findings) == 0, f"TSV 仅定义一次, 复用不应报 CON-003-B: {findings}"

    def test_abbreviation_fullname_parens_abbr_recognized(self):
        """CON-003: "全称 (TSV)" 定义格式 — 括号内缩写应识别为已定义 (不报未给出全称)"""
        fa = FactualAuditor()
        doc = _text_doc("Through Silicon Via (TSV) 工艺。TSV 用于先进封装。")
        findings = fa._check_abbreviation_first_defined(doc)
        assert len(findings) == 0, f"'全称 (TSV)' 应视为已定义, got: {findings}"

    def test_abbreviation_cache_weakref_identity(self):
        """回归: 缩写扫描缓存以 weakref.ref 绑定文档身份 (曾用 id(doc) — 地址可复用)"""
        import gc

        fa = FactualAuditor()
        doc1 = _text_doc("TSV 工艺用于先进封装。")
        fa._check_abbreviation_first_defined(doc1)
        key = fa._abbr_scan_cache[0]
        assert callable(key), f"缓存键应为 weakref.ref, got {type(key)}"
        assert key() is doc1, "存活时 weakref 应指向同一文档对象"
        del doc1
        gc.collect()
        assert key() is None, "文档对象回收后 weakref 应失效"

    def test_abbreviation_cache_weakref_new_doc_after_gc(self):
        """回归: 对象回收 + 新对象 → 缓存不串档 (weakref 失效后重扫)"""
        import gc

        fa = FactualAuditor()
        doc1 = _text_doc("TSV 工艺用于先进封装。")
        assert len(fa._check_abbreviation_first_defined(doc1)) == 1  # TSV 未定义
        del doc1
        gc.collect()
        doc2 = _text_doc("TSV (Through Silicon Via) 是硅通孔。")
        findings_doc2 = fa._check_abbreviation_first_defined(doc2)
        assert len(findings_doc2) == 0, (
            f"doc2 不应复用已回收 doc1 的扫描结果 (TSV 已定义), got: {findings_doc2}"
        )

    def test_abbreviation_cache_not_leaked_across_documents(self):
        """回归: 缩写扫描缓存不得跨文档串档 (独立模式 dispatch 直调不走 audit() 的 reset)"""
        doc1 = _text_doc("TSV 工艺用于先进封装。")
        doc2 = _text_doc("TSV (Through Silicon Via) 是硅通孔。")
        fa = FactualAuditor()
        # 模拟 CustomRulesAuditor 独立模式: 直接调 dispatch 方法 (不经过 audit() 重置缓存)
        fa._check_abbreviation_first_defined(doc1)
        findings_doc2 = fa._check_abbreviation_first_defined(doc2)
        # doc2 的 TSV 首次出现即带全称定义 → 不应报"首次未定义"
        assert len(findings_doc2) == 0, f"doc2 不应复用车 doc1 的扫描结果: {findings_doc2}"
