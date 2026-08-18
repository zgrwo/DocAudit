"""AutoFixer 单元测试 — 字体修复 / 间距修复 / 原子写入 / 链式修复"""

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from src.engines.autofix import AutoFixer


@pytest.fixture
def fixer():
    return AutoFixer(allowed_fonts=["微软雅黑", "Arial"])


@pytest.fixture
def sample_pptx(tmp_path):
    """创建含非标准字体的测试 PPTX"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "测试文本"
    run.font.name = "宋体"  # 非标准字体
    run.font.size = Pt(10)  # 过小字号
    path = tmp_path / "test_input.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def spacing_pptx(tmp_path):
    """创建含中英文紧排的测试 PPTX"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "使用FinFET技术"
    run.font.name = "Arial"
    path = tmp_path / "test_spacing.pptx"
    prs.save(str(path))
    return path


class TestFixPptx:
    """fix_pptx: 字体标准化 + 字号修正"""

    def test_nonstandard_font_replaced(self, fixer, sample_pptx, tmp_path):
        out = tmp_path / "output.pptx"
        fixer.fix_pptx(sample_pptx, out)
        assert out.exists()
        assert fixer.fix_count >= 1
        # 验证字体已替换
        prs = Presentation(str(out))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text:
                                assert run.font.name in fixer.allowed_fonts

    def test_small_font_size_corrected(self, fixer, sample_pptx, tmp_path):
        out = tmp_path / "output.pptx"
        fixer.fix_pptx(sample_pptx, out)
        prs = Presentation(str(out))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                assert run.font.size >= Pt(12)

    def test_atomic_write_no_corrupt_on_success(self, fixer, sample_pptx, tmp_path):
        """成功时目标文件存在且有效"""
        out = tmp_path / "output.pptx"
        result = fixer.fix_pptx(sample_pptx, out)
        assert result == out
        # 文件可被正常打开
        prs = Presentation(str(out))
        assert len(prs.slides) == 1

    def test_source_unchanged(self, fixer, sample_pptx, tmp_path):
        """源文件不被修改"""
        original_size = sample_pptx.stat().st_size
        out = tmp_path / "output.pptx"
        fixer.fix_pptx(sample_pptx, out)
        assert sample_pptx.stat().st_size == original_size


class TestFixSpacing:
    """fix_spacing: CJK-Latin 间距修复"""

    def test_spacing_inserted(self, fixer, spacing_pptx, tmp_path):
        out = tmp_path / "output_spacing.pptx"
        fixer.fix_spacing(spacing_pptx, out)
        assert out.exists()
        assert fixer.fix_count >= 1
        prs = Presentation(str(out))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if "FinFET" in text:
                        # 验证中英文之间有空格
                        assert " FinFET" in text or "FinFET " in text

    def test_unsupported_format_raises(self, fixer, tmp_path):
        """非 PPTX/DOCX 格式 → ValueError"""
        fake = tmp_path / "test.pdf"
        fake.write_text("fake")
        with pytest.raises(ValueError, match="仅支持"):
            fixer.fix_spacing(fake, tmp_path / "out.pdf")


class TestFixElementOverflow:
    """fix_element_overflow: 溢出元素移回边界"""

    def test_overflow_corrected(self, fixer, tmp_path):
        """负 left 值的 shape 被修正为 0"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1))
        # 手动设置为负值 (溢出左边界)
        txBox.left = -100000  # 负 EMU
        path = tmp_path / "overflow.pptx"
        prs.save(str(path))

        out = tmp_path / "overflow_fixed.pptx"
        fixer.fix_element_overflow(path, out)
        assert fixer.fix_count >= 1
        prs2 = Presentation(str(out))
        for shape in prs2.slides[0].shapes:
            assert shape.left >= 0

    def test_non_pptx_raises(self, fixer, tmp_path):
        fake = tmp_path / "test.docx"
        fake.write_text("fake")
        with pytest.raises(ValueError, match="仅支持 PPTX"):
            fixer.fix_element_overflow(fake, tmp_path / "out.docx")


class TestFixTitlePunctuation:
    """fix_title_punctuation: 标题末尾标点去除"""

    def test_trailing_punct_removed(self, fixer, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # title slide
        slide.shapes.title.text = "项目总结。"
        path = tmp_path / "title.pptx"
        prs.save(str(path))

        out = tmp_path / "title_fixed.pptx"
        fixer.fix_title_punctuation(path, out)
        assert fixer.fix_count >= 1
        prs2 = Presentation(str(out))
        title_text = prs2.slides[0].shapes.title.text
        assert not title_text.endswith("。")


class TestFixBulletStyle:
    """fix_bullet_style: 项目符号统一"""

    def test_dash_replaced_with_bullet(self, fixer, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        tf = txBox.text_frame
        tf.paragraphs[0].add_run().text = "- 第一项"
        p2 = tf.add_paragraph()
        p2.add_run().text = "- 第二项"
        path = tmp_path / "bullet.pptx"
        prs.save(str(path))

        out = tmp_path / "bullet_fixed.pptx"
        fixer.fix_bullet_style(path, out)
        assert fixer.fix_count >= 1
        prs2 = Presentation(str(out))
        for shape in prs2.slides[0].shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        assert not run.text.strip().startswith("- ")


class TestChainedFix:
    """链式修复: 前一步输出作为后一步输入"""

    def test_font_then_spacing(self, fixer, tmp_path):
        """字体修复 → 间距修复 链式执行"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tf = txBox.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = "使用FinFET技术"
        run.font.name = "宋体"
        path = tmp_path / "chain.pptx"
        prs.save(str(path))

        out = tmp_path / "chain_fixed.pptx"
        fixer.fix_pptx(path, out)
        font_fixes = fixer.fix_count
        fixer.fix_spacing(out, out)  # 链式: 在字体修复输出上继续
        spacing_fixes = fixer.fix_count

        assert font_fixes >= 1
        assert spacing_fixes >= 1
        assert out.exists()


class TestFixEastAsia:
    """eastAsia 中文字体修复: 字体替换必须同时写 latin 与 ea (FMT-001 中文修复生效)"""

    def test_fix_docx_replaces_east_asia(self, fixer, tmp_path):
        """DOCX: 非允许的 w:eastAsia 被替换为默认字体"""
        from docx import Document
        from docx.oxml.ns import qn

        doc = Document()
        p = doc.add_paragraph()
        run = p.add_run("中文正文")
        run.font.name = "Arial"  # 允许
        run._element.get_or_add_rPr().get_or_add_rFonts().set(qn("w:eastAsia"), "宋体")
        src = tmp_path / "ea_src.docx"
        doc.save(str(src))

        out = tmp_path / "ea_out.docx"
        fixer.fix_docx(src, out)
        assert out.exists()
        assert fixer.fix_count >= 1

        # 重开验证 eastAsia 已改为允许字体
        doc2 = Document(str(out))
        run2 = next(r for para in doc2.paragraphs for r in para.runs)
        rPr2 = run2._element.rPr
        assert rPr2 is not None and rPr2.rFonts is not None
        ea2 = rPr2.rFonts.get(qn("w:eastAsia"))
        assert ea2 in fixer.allowed_fonts, f"eastAsia 应变为允许字体, got {ea2}"
        # latin 保持不变 (Arial 本就在允许列表)
        assert run2.font.name == "Arial"

    def test_fix_docx_east_asia_idempotent(self, fixer, tmp_path):
        """DOCX: 二次运行幂等 (fix_count == 0)"""
        from docx import Document
        from docx.oxml.ns import qn

        doc = Document()
        p = doc.add_paragraph()
        run = p.add_run("中文正文")
        run.font.name = "Arial"
        run._element.get_or_add_rPr().get_or_add_rFonts().set(qn("w:eastAsia"), "宋体")
        src = tmp_path / "ea_src.docx"
        doc.save(str(src))

        out1 = tmp_path / "ea_1.docx"
        fixer.fix_docx(src, out1)
        assert fixer.fix_count >= 1
        out2 = tmp_path / "ea_2.docx"
        fixer.fix_docx(out1, out2)
        assert fixer.fix_count == 0, f"二次运行应幂等, fix_count={fixer.fix_count}"

    def test_fix_pptx_replaces_ea(self, fixer, tmp_path):
        """PPTX: 非允许的 a:ea typeface 被替换为默认字体"""
        from lxml import etree
        from pptx import Presentation
        from pptx.oxml.ns import qn
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = "中文正文"
        run.font.name = "Arial"
        ea = etree.SubElement(run.font._rPr, qn("a:ea"))
        ea.set("typeface", "宋体")
        src = tmp_path / "ea_src.pptx"
        prs.save(str(src))

        out = tmp_path / "ea_out.pptx"
        fixer.fix_pptx(src, out)
        assert out.exists()
        assert fixer.fix_count >= 1

        # 重开验证 a:ea 已改为允许字体
        prs2 = Presentation(str(out))
        run2 = next(
            r
            for slide in prs2.slides
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for r in para.runs
        )
        rPr2 = run2.font._rPr
        ea2 = rPr2.find(qn("a:ea"))
        assert ea2 is not None
        assert ea2.get("typeface") in fixer.allowed_fonts, (
            f"a:ea 应变为允许字体, got {ea2.get('typeface')}"
        )
        assert run2.font.name == "Arial"  # latin 不变
