"""CLI 退出码契约测试 — 处理失败必须 exit 1 (README 声明, markdownlint 风格)。

背景: README 声明"处理失败 → 退出码 1"，但 cli.py 的 except Exception 分支只打印
错误不置失败标记，处理失败后仍按 findings 统计退出 (损坏文件实测 exit 0)。
本测试固化契约: 处理失败 → exit 1。
"""

import subprocess
import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent


def _run_cli(*args: str) -> subprocess.CompletedProcess:
    """在项目根目录运行 src/cli.py，返回子进程结果。"""
    return subprocess.run(
        [sys.executable, str(PROJECT_ROOT / "src" / "cli.py"), *args],
        capture_output=True,
        text=True,
        timeout=300,
        cwd=str(PROJECT_ROOT),
        encoding="utf-8",
        errors="replace",  # CLI 在 Windows 上输出 UTF-8
    )


def _write_clean_pptx(path: Path) -> None:
    """生成无 Error 级发现的 PPTX (含 CON-002 必含章节 + CON-004 每页结论)。

    依赖清单 (rules.md 当前阈值, 改动后需同步本函数):
    - CON-002 必含章节: 概述 / 工艺流程 / 关键参数 / 结论
    - CON-004 每页结论: 每页正文须含"结论"等关键词
    - STR-001 标题版式: slide_layouts[0] (标题版式)
    - FMT-001 允许字体: 占位符无显式字体 (默认字体)
    若测试失败, 先确认 rules.md 阈值变化; 断言失败信息会输出实际 findings 便于定位。
    """
    from pptx import Presentation

    prs = Presentation()
    for title in ("概述", "工艺流程", "关键参数", "结论"):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:
                ph.text = f"{title}内容。结论：本节符合预期，建议保留。"
    prs.save(str(path))


class TestCliExitCodes:
    """CLI 退出码契约。"""

    def test_corrupted_pptx_exits_1(self, tmp_path):
        """损坏文件 (非 PPTX 内容命名为 .pptx) → 处理失败 → exit 1。"""
        broken = tmp_path / "broken.pptx"
        broken.write_text("这不是一个真正的 PPTX 文件", encoding="utf-8")
        result = _run_cli(str(broken), "--rules", str(PROJECT_ROOT / "rules.md"))
        assert result.returncode == 1, (
            f"处理失败应 exit 1，实际 {result.returncode}\n{result.stdout[-500:]}"
        )
        assert "failed to process" in result.stdout

    def test_clean_file_exits_0(self, tmp_path):
        """正常文件 (无 Error 级发现) → exit 0。"""
        clean = tmp_path / "clean.pptx"
        _write_clean_pptx(clean)
        result = _run_cli(str(clean), "--rules", str(PROJECT_ROOT / "rules.md"))
        # 失败时输出完整 stdout (含 [SUMMARY] 计数与 Errors 列表) 便于对照 rules.md 阈值定位
        assert result.returncode == 0, (
            f"无 Error 应 exit 0，实际 {result.returncode}\n{result.stdout}"
        )

    def test_sample_pptx_with_errors_exits_1(self):
        """正常文件但存在 Error 级发现 (sample.pptx) → exit 1。"""
        sample = PROJECT_ROOT / "tests" / "fixtures" / "sample.pptx"
        if not sample.exists():
            import pytest

            pytest.skip("测试夹具 tests/fixtures/sample.pptx 不存在")
        result = _run_cli(str(sample), "--rules", str(PROJECT_ROOT / "rules.md"))
        assert result.returncode == 1, (
            f"有 Error 应 exit 1，实际 {result.returncode}\n{result.stdout[-500:]}"
        )
        assert "error(s) found" in result.stdout

    def test_nonexistent_path_exits_1(self):
        """路径不存在 → exit 1。"""
        missing = PROJECT_ROOT / "tests" / "fixtures" / "no_such_file.pptx"
        result = _run_cli(str(missing), "--rules", str(PROJECT_ROOT / "rules.md"))
        assert result.returncode == 1, (
            f"路径不存在应 exit 1，实际 {result.returncode}\n{result.stdout[-500:]}"
        )


class TestCliOutputContract:
    """-o 导出契约 (M6/M9): 导出失败 → exit 1；不支持的扩展名 → exit 2。"""

    def test_output_dir_missing_exits_1(self, tmp_path):
        """-o 指向不存在的目录 → 报告未生成 → exit 1 且输出含错误。"""
        clean = tmp_path / "clean.pptx"
        _write_clean_pptx(clean)
        out = tmp_path / "no_such_dir" / "report.html"
        result = _run_cli(str(clean), "--rules", str(PROJECT_ROOT / "rules.md"), "-o", str(out))
        assert result.returncode == 1, (
            f"导出失败应 exit 1，实际 {result.returncode}\n{result.stdout[-500:]}"
        )
        assert "导出失败" in result.stdout or "Processing failed" in result.stdout

    def test_output_unsupported_extension_exits_2(self, tmp_path):
        """-o 非 .html/.json 扩展名 → 打印错误 + exit 2 (用法错误)。"""
        clean = tmp_path / "clean.pptx"
        _write_clean_pptx(clean)
        out = tmp_path / "report.txt"
        result = _run_cli(str(clean), "--rules", str(PROJECT_ROOT / "rules.md"), "-o", str(out))
        assert result.returncode == 2, (
            f"不支持扩展名应 exit 2，实际 {result.returncode}\n{result.stdout[-500:]}"
        )
        assert "不支持的输出格式" in result.stdout

    def test_fix_unsupported_format_prints_hint(self, tmp_path):
        """--fix 对 MD → 打印不支持自动修复提示 (L7, 不改变行为)。"""
        md = tmp_path / "doc.md"
        md.write_text("# 概述\n\n测试内容。\n\n## 结论\n\n测试完成。\n", encoding="utf-8")
        result = _run_cli(str(md), "--rules", str(PROJECT_ROOT / "rules.md"), "--fix")
        assert "暂不支持自动修复" in result.stdout
