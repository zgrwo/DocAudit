"""Converter 单元测试 — DocxConverter 字段正确性 + MarkdownConverter 边界"""

import io
import logging
from pathlib import Path

import pytest
from docx import Document as DocxDocument
from docx.shared import Pt

from src.converters.docx_converter import DocxConverter
from src.converters.md_converter import MarkdownConverter
from src.converters.pptx_converter import PptxConverter
from src.models.document import Page, PageElement


@pytest.fixture
def docx_converter():
    return DocxConverter()


@pytest.fixture
def md_converter():
    return MarkdownConverter()


@pytest.fixture
def sample_docx(tmp_path):
    """创建含标题+正文+表格的测试 DOCX"""
    doc = DocxDocument()
    # 标题
    doc.add_heading("测试标题", level=1)
    # 正文段落
    p = doc.add_paragraph("这是正文内容。")
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(12)
    # 二级标题
    doc.add_heading("第二节", level=2)
    doc.add_paragraph("第二节的内容。")
    # 表格
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A1"
    table.cell(0, 1).text = "B1"
    table.cell(1, 0).text = "A2"
    table.cell(1, 1).text = "B2"

    path = tmp_path / "test.docx"
    doc.save(str(path))
    return path


class TestDocxConverter:
    """DocxConverter 转换正确性"""

    def test_can_handle(self, docx_converter):
        assert docx_converter.can_handle("test.docx")
        assert docx_converter.can_handle("test.doc")
        assert not docx_converter.can_handle("test.pptx")
        assert not docx_converter.can_handle("test.pdf")

    def test_convert_returns_document(self, docx_converter, sample_docx):
        doc = docx_converter.convert(str(sample_docx))
        assert doc.format == "docx"
        assert doc.source_path == str(sample_docx)
        assert len(doc.pages) >= 1

    def test_heading_detected_as_title(self, docx_converter, sample_docx):
        """标题样式段落被识别（style_name 含 Heading）且 is_title=True"""
        doc = docx_converter.convert(str(sample_docx))
        # python-docx add_heading 不设置 outlineLvl，但样式名含 Heading
        all_elements = [e for p in doc.pages for e in p.flattened_elements]
        heading_elements = [
            e for e in all_elements if e.style_name and "heading" in e.style_name.lower()
        ]
        assert len(heading_elements) >= 1
        # HIGH-1: 样式型标题（无段落级 outlineLvl）也必须判为标题
        for e in heading_elements:
            assert e.is_title is True, f"标题元素应 is_title=True: {e.style_name!r}"

    def test_heading_level_from_style_fallback(self, docx_converter, tmp_path):
        """HIGH-1: add_heading 生成的标题（无段落级 outlineLvl，样式级定义）也能提取 level

        样式级回退顺序: ① style 的 w:pPr/w:outlineLvl → ② 样式名 "Heading N"→N-1
        """
        doc = DocxDocument()
        doc.add_heading("一级标题", level=1)
        doc.add_heading("二级标题", level=2)
        path = tmp_path / "heading_style.docx"
        doc.save(str(path))

        result = docx_converter.convert(str(path))
        elems = [e for p in result.pages for e in p.flattened_elements if e.paragraphs]
        h1 = next(e for e in elems if e.paragraphs[0].text == "一级标题")
        assert h1.paragraphs[0].level == 0  # "Heading 1" → 0-based level 0
        assert h1.is_title is True
        h2 = next(e for e in elems if e.paragraphs[0].text == "二级标题")
        assert h2.paragraphs[0].level == 1  # "Heading 2" → 0-based level 1
        assert h2.is_title is True  # level 1 <= 1 → 仍为标题
        # 语义分页: 两个 H1/H2 级标题均为页面边界 (旧实现 level=None → 整篇单页)
        assert len(result.pages) >= 2, f"样式型标题应触发分页, got {len(result.pages)} 页"
        assert any(e.paragraphs[0].text == "一级标题" for e in result.pages[0].flattened_elements)

    def test_heading_level_from_style_name_regex(self, docx_converter, tmp_path):
        """HIGH-1: 样式无 outlineLvl 时按样式名 "Heading N" 回退 (N-1)"""
        from docx.oxml.ns import qn

        doc = DocxDocument()
        p = doc.add_paragraph("手写标题")
        p.style = doc.styles["Heading 3"]
        # 移除样式级 outlineLvl，模拟无大纲级别的 Heading 样式
        style_el = p.style._element
        pPr = style_el.find(qn("w:pPr"))
        if pPr is not None:
            ol = pPr.find(qn("w:outlineLvl"))
            if ol is not None:
                pPr.remove(ol)
        path = tmp_path / "heading_regex.docx"
        doc.save(str(path))

        result = docx_converter.convert(str(path))
        elems = [e for p in result.pages for e in p.flattened_elements if e.paragraphs]
        h3 = next(e for e in elems if e.paragraphs[0].text == "手写标题")
        assert h3.paragraphs[0].level == 2  # "Heading 3" → 2

    def test_style_name_moved_from_shape_name(self, docx_converter, sample_docx):
        """DOCX 段落样式名移到 style_name 字段, shape_name 恢复 None"""
        doc = docx_converter.convert(str(sample_docx))
        all_elements = [e for p in doc.pages for e in p.flattened_elements]
        heading_elems = [
            e for e in all_elements if e.style_name and "heading" in e.style_name.lower()
        ]
        assert len(heading_elems) >= 1, "Heading 样式段落应有 style_name"
        # shape_name 不再被 DOCX 样式名占用 (保留给 PPTX shape 名)
        for e in all_elements:
            assert e.shape_name is None, f"DOCX 元素 shape_name 应为 None, got {e.shape_name!r}"

    def test_paragraph_text_preserved(self, docx_converter, sample_docx):
        """正文文本完整保留"""
        doc = docx_converter.convert(str(sample_docx))
        assert "这是正文内容。" in doc.all_text
        assert "第二节的内容。" in doc.all_text

    def test_font_info_extracted(self, docx_converter, sample_docx):
        """Run 级字体信息被提取"""
        doc = docx_converter.convert(str(sample_docx))
        found_arial = False
        for page in doc.pages:
            for elem in page.flattened_elements:
                for para in elem.paragraphs:
                    for run in para.runs:
                        if run.font_name == "Arial":
                            found_arial = True
                            assert run.font_size == 12.0
        assert found_arial, "Should find Arial font in converted document"

    def test_table_converted(self, docx_converter, sample_docx):
        """表格被正确转换"""
        doc = docx_converter.convert(str(sample_docx))
        table_elements = [e for p in doc.pages for e in p.flattened_elements if e.type == "table"]
        assert len(table_elements) >= 1
        # 验证表格内容
        table_elem = table_elements[0]
        all_cell_text = [cell.text for row in table_elem.tables for cell in row]
        assert "A1" in all_cell_text
        assert "B2" in all_cell_text

    def test_metadata_extracted(self, docx_converter, sample_docx):
        """元数据（字数）被提取"""
        doc = docx_converter.convert(str(sample_docx))
        assert doc.metadata.word_count is not None
        assert doc.metadata.word_count > 0

    def test_table_cell_colors_extracted(self, docx_converter, tmp_path):
        """表格单元格: w:shd 底色 + 字体色提取 (FMT-008 数据源)"""
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        from docx.shared import RGBColor

        doc = DocxDocument()
        table = doc.add_table(rows=2, cols=2)
        # 深蓝底 + 白字
        cell = table.cell(0, 0)
        cell.text = "深底浅字"
        tcPr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:fill"), "1E3A5F")
        tcPr.append(shd)
        cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # 无底色单元格
        table.cell(0, 1).text = "无底色"

        path = tmp_path / "table_colors.docx"
        doc.save(str(path))

        result = docx_converter.convert(str(path))
        tables = [e for p in result.pages for e in p.flattened_elements if e.type == "table"]
        assert tables, "应有表格元素"
        cells = [c for row in tables[0].tables for c in row]
        c00 = next(c for c in cells if (c.row, c.col) == (0, 0))
        assert c00.fill_color == "1E3A5F"
        assert c00.font_color == "FFFFFF"
        c01 = next(c for c in cells if (c.row, c.col) == (0, 1))
        assert c01.fill_color is None  # 无底色 → None (不误报)

    def test_empty_docx_no_crash(self, docx_converter, tmp_path):
        """空 DOCX 不崩溃"""
        doc = DocxDocument()
        path = tmp_path / "empty.docx"
        doc.save(str(path))
        result = docx_converter.convert(str(path))
        assert result.format == "docx"
        assert len(result.pages) >= 1

    def test_docx_run_font_name_east_asia_extracted(self, docx_converter, tmp_path):
        """eastAsia 中文字体被提取到 Run.font_name_east_asia（FMT-001 中文判定数据源）"""
        from docx.oxml.ns import qn

        doc = DocxDocument()
        p = doc.add_paragraph()
        run = p.add_run("中文正文")
        rPr = run._element.get_or_add_rPr()
        rPr.get_or_add_rFonts().set(qn("w:eastAsia"), "宋体")
        rPr.rFonts.set(qn("w:ascii"), "Arial")
        path = tmp_path / "ea.docx"
        doc.save(str(path))

        result = docx_converter.convert(str(path))
        run_model = next(
            r
            for page in result.pages
            for e in page.flattened_elements
            for para in e.paragraphs
            for r in para.runs
        )
        assert run_model.font_name_east_asia == "宋体"
        assert run_model.font_name == "Arial"  # latin 字体不受影响

    def test_docx_run_font_name_east_asia_none_when_missing(self, docx_converter, tmp_path):
        """无 eastAsia 字体时 font_name_east_asia 为 None"""
        doc = DocxDocument()
        p = doc.add_paragraph()
        p.add_run("纯文本")
        path = tmp_path / "plain.docx"
        doc.save(str(path))

        result = docx_converter.convert(str(path))
        run_model = next(
            r
            for page in result.pages
            for e in page.flattened_elements
            for para in e.paragraphs
            for r in para.runs
        )
        assert run_model.font_name_east_asia is None

    def test_single_run_exception_keeps_other_runs(self, docx_converter, tmp_path, monkeypatch):
        """F9: 单个 run 格式提取失败只跳过该 run，保留其余 run 的格式

        (原实现外层 except 包整个 runs 循环，一个 run 异常会丢掉整段格式)
        """
        import docx.text.run as docx_run_module

        doc = DocxDocument()
        p = doc.add_paragraph()
        p.add_run("正常一")
        p.add_run("异常二")
        p.add_run("正常三")
        for r in p.runs:
            r.font.name = "Arial"
        path = tmp_path / "flaky.docx"
        doc.save(str(path))

        real_font = docx_run_module.Run.font

        def flaky_font(self):
            if self.text == "异常二":
                raise ValueError("模拟单个 run 格式损坏")
            return real_font.fget(self)

        monkeypatch.setattr(docx_run_module.Run, "font", property(flaky_font))

        result = docx_converter.convert(str(path))
        run_models = [
            r
            for page in result.pages
            for e in page.flattened_elements
            for para in e.paragraphs
            for r in para.runs
        ]
        # 异常 run 被跳过，其余 run 完整保留且格式不丢
        assert [r.text for r in run_models] == ["正常一", "正常三"]
        for r in run_models:
            assert r.font_name == "Arial", f"run {r.text!r} 格式应保留"

    def test_line_spacing_length_converted_to_pt(self, docx_converter, tmp_path):
        """F5: DOCX 固定行距 (Twips Length) 转为 pt 数值，满足 Paragraph.line_spacing float|None 契约"""
        from docx.shared import Pt

        doc = DocxDocument()
        p = doc.add_paragraph("固定行距段落")
        p.paragraph_format.line_spacing = Pt(24)
        path = tmp_path / "spacing.docx"
        doc.save(str(path))

        result = docx_converter.convert(str(path))
        para = next(
            pp
            for page in result.pages
            for e in page.flattened_elements
            for pp in e.paragraphs
            if pp.text.strip()
        )
        assert para.line_spacing == 24.0
        assert isinstance(para.line_spacing, float), (
            f"line_spacing 应为 float, got {type(para.line_spacing).__name__}"
        )


class TestMarkdownConverter:
    """MarkdownConverter 边界测试"""

    def test_can_handle(self, md_converter):
        assert md_converter.can_handle("test.md")
        assert md_converter.can_handle("test.markdown")
        assert md_converter.can_handle("test.txt")
        assert not md_converter.can_handle("test.docx")

    def test_heading_levels(self, md_converter, tmp_path):
        """Markdown 标题层级正确映射"""
        md = "# H1\n\n## H2\n\n### H3\n\n正文\n"
        path = tmp_path / "test.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        # 应有标题元素
        all_paras = doc.all_paragraphs
        levels = [p.level for p in all_paras if p.level is not None]
        assert 1 in levels
        assert 2 in levels

    def test_code_block_preserved(self, md_converter, tmp_path):
        """围栏代码块完整保留"""
        md = "# Title\n\n```python\nprint('hello')\n```\n\nAfter code.\n"
        path = tmp_path / "code.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        assert "print('hello')" in doc.all_text

    def test_table_parsed(self, md_converter, tmp_path):
        """GFM 表格被解析为 table 元素"""
        md = "# Data\n\n| Name | Age |\n|------|-----|\n| Alice | 30 |\n"
        path = tmp_path / "table.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        table_elems = [e for p in doc.pages for e in p.flattened_elements if e.type == "table"]
        assert len(table_elems) >= 1

    def test_list_item_with_pipe_not_table(self, md_converter, tmp_path):
        """回归: 列表项含竖线不得误判为表格 (修复: 曾把 '- 项 A | 内容' 解析成 table)"""
        md = "- 列表项 A | 内容\n- 普通列表项\n"
        path = tmp_path / "list.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        table_elems = [e for p in doc.pages for e in p.flattened_elements if e.type == "table"]
        assert len(table_elems) == 0, f"列表项不应被解析为表格: {table_elems}"
        # 两个列表项都应在文本中完整保留
        assert "- 列表项 A | 内容" in doc.all_text
        assert "- 普通列表项" in doc.all_text

    def test_gbk_encoding_fallback(self, md_converter, tmp_path):
        """GBK 编码文件正确读取"""
        path = tmp_path / "gbk.md"
        path.write_text("# 标题\n\n中文内容\n", encoding="gbk")
        doc = md_converter.convert(str(path))
        assert "中文内容" in doc.all_text

    def test_empty_file_no_crash(self, md_converter, tmp_path):
        """空文件不崩溃"""
        path = tmp_path / "empty.md"
        path.write_text("", encoding="utf-8")
        doc = md_converter.convert(str(path))
        assert doc.format == "md"

    def test_page_index_continuous_with_empty_chunks(self, md_converter, tmp_path):
        """F10: 空 chunk 被跳过时页面 index/slide_number 保持连续 (旧实现出现 0,3,5 空洞)"""
        md = "# 第一页\n\n# 第二页\n\n\n\n# 第三页\n"
        path = tmp_path / "pages.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        assert [p.index for p in doc.pages] == [0, 1, 2], (
            f"页面 index 应连续, got {[p.index for p in doc.pages]}"
        )
        assert [p.slide_number for p in doc.pages] == [1, 2, 3], (
            f"slide_number 应连续, got {[p.slide_number for p in doc.pages]}"
        )

    def test_utf8_bom_frontmatter_still_detected(self, md_converter, tmp_path):
        """UTF-8 BOM 文件的 frontmatter 仍能正确解析 (BOM 剥离)"""
        path = tmp_path / "bom.md"
        content = "---\ntitle: BOM 文档\nauthor: tester\n---\n\n# 标题\n\n内容\n"
        path.write_bytes(b"\xef\xbb\xbf" + content.encode("utf-8"))
        doc = md_converter.convert(str(path))
        assert doc.metadata.title == "BOM 文档"
        assert doc.metadata.author == "tester"
        assert "\ufeff" not in doc.all_text


def _inject_fake_docling(monkeypatch, converter_cls) -> None:
    """封闭式注入假 docling 模块到 sys.modules (不触发真实 import, CI 无依赖也可跑)。"""
    import sys
    import types

    fake_pkg = types.ModuleType("docling")
    fake_conv = types.ModuleType("docling.document_converter")
    fake_conv.DocumentConverter = converter_cls
    fake_pkg.document_converter = fake_conv
    monkeypatch.setitem(sys.modules, "docling", fake_pkg)
    monkeypatch.setitem(sys.modules, "docling.document_converter", fake_conv)


class TestPdfConverter:
    """PdfConverter 基本行为 (docling 可选依赖, 测试用封闭式 mock 不依赖安装)"""

    def test_can_handle(self):
        from src.converters.pdf_converter import PdfConverter

        cvt = PdfConverter()
        assert cvt.can_handle("test.pdf")
        assert cvt.can_handle("TEST.PDF")
        assert not cvt.can_handle("test.docx")

    def test_missing_docling_raises_helpful_error(self, tmp_path):
        """docling 未安装 → 抛出含安装指引的 ImportError"""
        pytest.importorskip("src.converters.pdf_converter")
        try:
            import docling  # noqa: F401

            pytest.skip("docling 已安装，无法验证缺失依赖路径")
        except ImportError:
            pass
        from src.converters.pdf_converter import PdfConverter

        fake = tmp_path / "fake.pdf"
        fake.write_bytes(b"%PDF-1.4 fake content")
        with pytest.raises(ImportError, match="docling"):
            PdfConverter().convert(str(fake))

    def test_convert_positive_path_with_cells_api(self, tmp_path, monkeypatch):
        """正路径: cells API 页面 → 文本完整保留 (docling mock, 封闭式注入不依赖安装)"""
        from src.converters.pdf_converter import PdfConverter

        class FakeCell:
            def __init__(self, text, row=0, col=0):
                self.text = text
                self.row = row
                self.col = col

        class FakePage:
            def __init__(self):
                self.cells = [FakeCell("第一段文本", 0, 0)]

        class FakeDoclingDoc:
            pages = {1: FakePage()}

        class FakeResult:
            document = FakeDoclingDoc()

        class FakeDoclingConverter:
            def __init__(self, *a, **kw):
                pass

            def convert(self, path):
                return FakeResult()

        # 封闭式注入: 直接挂 sys.modules, 不 import 真实 docling (CI 无该依赖也通过)
        _inject_fake_docling(monkeypatch, FakeDoclingConverter)
        fake = tmp_path / "sample.pdf"
        fake.write_bytes(b"%PDF-1.4 fake")

        doc = PdfConverter().convert(str(fake))
        assert doc.format == "pdf"
        assert doc.metadata.page_count == 1
        assert "第一段文本" in doc.all_text

    def test_convert_positive_path_with_tables_api(self, tmp_path, monkeypatch):
        """正路径: tables API + export_to_dataframe → table 元素 (docling/pandas mock 封闭)"""
        from src.converters.pdf_converter import PdfConverter

        class FakeDataFrame:
            """最小 iterrows 假实现, 避免依赖真实 pandas"""

            def iterrows(self):
                yield (0, ["表头A", "表头B"])
                yield (1, ["值1", "值2"])

        class FakeTable:
            def export_to_dataframe(self):
                return FakeDataFrame()

        class FakePage:
            def __init__(self):
                self.tables = [FakeTable()]

        class FakeDoclingDoc:
            pages = {1: FakePage()}

        class FakeResult:
            document = FakeDoclingDoc()

        class FakeDoclingConverter:
            def __init__(self, *a, **kw):
                pass

            def convert(self, path):
                return FakeResult()

        _inject_fake_docling(monkeypatch, FakeDoclingConverter)
        # 假 pandas 模块: 让 _convert_docling_table 的 import pandas 检查通过
        import sys
        import types

        monkeypatch.setitem(sys.modules, "pandas", types.ModuleType("pandas"))

        fake = tmp_path / "table.pdf"
        fake.write_bytes(b"%PDF-1.4 fake")

        doc = PdfConverter().convert(str(fake))
        tables = [e for p in doc.pages for e in p.flattened_elements if e.type == "table"]
        assert len(tables) == 1, f"期望 1 个表格元素, got {len(tables)}"
        cells = [c.text for row in tables[0].tables for c in row]
        assert cells == ["表头A", "表头B", "值1", "值2"]


class TestPptxConverterTableColors:
    """PPTX 表格单元格底色/字体色提取 (FMT-008 数据源)"""

    def test_table_cell_colors_extracted(self, tmp_path):
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches

        from src.converters.pptx_converter import PptxConverter

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1))
        table = shape.table
        # 深蓝底 + 白字
        cell = table.cell(0, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
        cell.text = "深底浅字"
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # 无填充单元格
        table.cell(0, 1).text = "无填充"

        path = tmp_path / "table_colors.pptx"
        prs.save(str(path))

        doc = PptxConverter().convert(str(path))
        tables = [e for p in doc.pages for e in p.flattened_elements if e.type == "table"]
        assert tables, "应有表格元素"
        cells = [c for row in tables[0].tables for c in row]
        c00 = next(c for c in cells if (c.row, c.col) == (0, 0))
        assert c00.fill_color == "1E3A5F"
        assert c00.font_color == "FFFFFF"
        c01 = next(c for c in cells if (c.row, c.col) == (0, 1))
        assert c01.fill_color is None  # 无填充 → None (不误报)


class TestPptxConverterLineSpacing:
    """PPTX 段落行距类型契约 (F5)"""

    def test_line_spacing_length_converted_to_pt(self, tmp_path):
        """F5: PPTX 固定行距 (Centipoints Length) 转为 pt 数值，满足 float|None 契约"""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        from src.converters.pptx_converter import PptxConverter

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        para = txBox.text_frame.paragraphs[0]
        para.text = "固定行距段落"
        para.line_spacing = Pt(28)
        path = tmp_path / "spacing.pptx"
        prs.save(str(path))

        doc = PptxConverter().convert(str(path))
        para_model = next(
            pp
            for page in doc.pages
            for e in page.flattened_elements
            for pp in e.paragraphs
            if pp.text.strip()
        )
        assert para_model.line_spacing == 28.0
        assert isinstance(para_model.line_spacing, float), (
            f"line_spacing 应为 float, got {type(para_model.line_spacing).__name__}"
        )


class TestPptxConverterRprReadOnly:
    """F4: run 属性访问不得变异 XML (font._rPr 会 get_or_add 创建空 a:rPr)"""

    def test_ea_lookup_read_only_no_rpr_created(self):
        """无 a:rPr 的 run: _ea_typeface 只读查找返回 None 且不创建 a:rPr"""
        from pptx import Presentation
        from pptx.oxml.ns import qn
        from pptx.util import Inches

        from src.converters.pptx_converter import _ea_typeface

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = "无格式文本"
        assert run._r.find(qn("a:rPr")) is None  # 前置: 无 rPr

        assert _ea_typeface(run) is None  # 不存在 → None
        assert run._r.find(qn("a:rPr")) is None, (
            "只读查找不得凭空创建 a:rPr (font._rPr 会 get_or_add)"
        )

    def test_ea_lookup_reads_existing_rpr(self):
        """已有 a:rPr/a:ea 的 run: _ea_typeface 能读到 typeface"""
        from lxml import etree
        from pptx import Presentation
        from pptx.oxml.ns import qn
        from pptx.util import Inches

        from src.converters.pptx_converter import _ea_typeface

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = "中文正文"
        rPr = run.font._rPr  # 测试侧显式创建 rPr (python-pptx 无 get_or_add_ea)
        ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", "宋体")

        assert _ea_typeface(run) == "宋体"


class TestPptxConverterEastAsia:
    """PPTX a:ea (eastAsia) 中文字体提取"""

    def test_run_font_name_east_asia_extracted(self, tmp_path):
        """a:ea typeface 被提取到 Run.font_name_east_asia"""
        from lxml import etree
        from pptx import Presentation
        from pptx.oxml.ns import qn
        from pptx.util import Inches

        from src.converters.pptx_converter import PptxConverter

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = "中文正文"
        run.font.name = "Arial"
        # 手工创建 a:ea (python-pptx 1.0.2 无 get_or_add_ea)
        rPr = run.font._rPr
        ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", "宋体")

        path = tmp_path / "ea.pptx"
        prs.save(str(path))

        doc = PptxConverter().convert(str(path))
        run_model = next(
            r
            for page in doc.pages
            for e in page.flattened_elements
            for para in e.paragraphs
            for r in para.runs
        )
        assert run_model.font_name_east_asia == "宋体"
        assert run_model.font_name == "Arial"  # latin 字体不受影响

    def test_run_font_name_east_asia_none_when_missing(self, tmp_path):
        """无 a:ea 时 font_name_east_asia 为 None"""
        from pptx import Presentation
        from pptx.util import Inches

        from src.converters.pptx_converter import PptxConverter

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = "纯英文 text"
        run.font.name = "Arial"
        path = tmp_path / "plain.pptx"
        prs.save(str(path))

        doc = PptxConverter().convert(str(path))
        run_model = next(
            r
            for page in doc.pages
            for e in page.flattened_elements
            for para in e.paragraphs
            for r in para.runs
        )
        assert run_model.font_name_east_asia is None


class TestDocxNestedContentWarnings:
    """DOCX 未解析嵌套内容 (文本框/页眉页脚/脚注) 跳过时输出 logger.warning"""

    @staticmethod
    def _make_textbox_docx(path) -> None:
        """构造含 w:txbxContent 文本框的 docx (python-docx 无文本框 API，手工注入 XML)"""
        from docx.oxml import parse_xml
        from docx.oxml.ns import nsdecls

        WPS = 'xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape"'
        doc = DocxDocument()
        p = doc.add_paragraph("正文")
        run = p.add_run()
        drawing_xml = (
            "<w:drawing {w}>"
            "<wp:inline {wp}>"
            '<a:graphic {a}><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/picture">'
            "<wps:wsp {wps}><wps:txbx><w:txbxContent>"
            "<w:p><w:r><w:t>文本框内容</w:t></w:r></w:p>"
            "</w:txbxContent></wps:txbx></wps:wsp>"
            "</a:graphicData></a:graphic></wp:inline></w:drawing>"
        ).format(w=nsdecls("w"), wp=nsdecls("wp"), a=nsdecls("a"), wps=WPS)
        run._element.append(parse_xml(drawing_xml))
        doc.save(str(path))

    def test_textbox_content_skip_warned(self, docx_converter, tmp_path, caplog):
        """含 w:txbxContent 文本框 → logger.warning 提示跳过范围"""
        path = tmp_path / "txbx.docx"
        self._make_textbox_docx(path)
        with caplog.at_level(logging.WARNING, logger="src.converters.docx_converter"):
            docx_converter.convert(str(path))
        assert any("文本框" in r.message for r in caplog.records), (
            f"应输出文本框跳过警告, got: {[r.message for r in caplog.records]}"
        )

    def test_header_footer_skip_warned(self, docx_converter, tmp_path, caplog):
        """含页眉内容 → logger.warning 提示跳过范围"""
        path = tmp_path / "hdr.docx"
        doc = DocxDocument()
        doc.sections[0].header.paragraphs[0].text = "页眉内容"
        doc.save(str(path))
        with caplog.at_level(logging.WARNING, logger="src.converters.docx_converter"):
            docx_converter.convert(str(path))
        assert any("页眉" in r.message for r in caplog.records), (
            f"应输出页眉/页脚跳过警告, got: {[r.message for r in caplog.records]}"
        )

    def test_footnote_skip_warned(self, docx_converter, tmp_path, caplog):
        """含脚注部件 → logger.warning 提示跳过范围"""
        from docx.opc.constants import CONTENT_TYPE, RELATIONSHIP_TYPE
        from docx.opc.packuri import PackURI
        from docx.opc.part import Part

        path = tmp_path / "fn.docx"
        doc = DocxDocument()
        fn_xml = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            '<w:footnote w:type="separator" w:id="-1"><w:p><w:r><w:separator/></w:r></w:p></w:footnote>'
            '<w:footnote w:id="1"><w:p><w:r><w:t xml:space="preserve"> 脚注内容</w:t></w:r></w:p></w:footnote>'
            "</w:footnotes>"
        )
        part = Part(
            PackURI("/word/footnotes.xml"),
            CONTENT_TYPE.WML_FOOTNOTES,
            fn_xml.encode("utf-8"),
            doc.part.package,
        )
        doc.part.relate_to(part, RELATIONSHIP_TYPE.FOOTNOTES)
        doc.save(str(path))
        with caplog.at_level(logging.WARNING, logger="src.converters.docx_converter"):
            docx_converter.convert(str(path))
        assert any("脚注" in r.message for r in caplog.records), (
            f"应输出脚注跳过警告, got: {[r.message for r in caplog.records]}"
        )


class TestPptxConverterMemory:
    """P1-5: PPTX 图片/内嵌 Excel 不再整包载入内存 (WP-D)"""

    # 1x1 透明 PNG — 最小合法图片文件
    MINIMAL_PNG = bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
        "0000000d49444154789c6360000002000154a24f9f0000000049454e44ae426082"
    )

    @staticmethod
    def _make_image_pptx(path) -> None:
        """构造含一张 PNG 图片的 PPTX"""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            io.BytesIO(TestPptxConverterMemory.MINIMAL_PNG),
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
        )
        prs.save(str(path))

    @staticmethod
    def _make_chart_pptx(path) -> None:
        """构造含内嵌 Excel 图表的 PPTX (python-pptx add_chart 会自动嵌入 xlsx part)"""
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        chart_data = CategoryChartData()
        chart_data.categories = ["A", "B", "C"]
        chart_data.add_series("系列1", (1, 2, 3))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), chart_data
        )
        prs.save(str(path))

    def test_image_element_keeps_ext_only(self, tmp_path):
        """图片元素只保留 image_ext，不再暴露 image_blob 字段"""
        path = tmp_path / "img.pptx"
        self._make_image_pptx(path)
        doc = PptxConverter().convert(str(path))
        images = [e for p in doc.pages for e in p.flattened_elements if e.type == "image"]
        assert images, "应有图片元素"
        img = images[0]
        assert img.image_ext == "png"
        assert not hasattr(img, "image_blob"), "image_blob 字段应已移除 (整图内存红线)"

    def test_pageelement_model_has_no_image_blob(self):
        """模型层: PageElement / Page 均不再有 image_blob 字段"""
        elem = PageElement(type="image")
        assert not hasattr(elem, "image_blob")
        page = Page(index=0)
        assert not hasattr(page, "image_blob")

    def test_chart_keeps_type_but_no_data_blob(self, tmp_path):
        """图表保留 chart_type；chart_data 字段已移除 (无消费者, F7)"""
        path = tmp_path / "chart.pptx"
        self._make_chart_pptx(path)
        doc = PptxConverter().convert(str(path))
        charts = [e for p in doc.pages for e in p.flattened_elements if e.type == "chart"]
        assert charts, "应有图表元素"
        chart = charts[0]
        assert chart.chart_type is not None
        assert not hasattr(chart, "chart_data"), "chart_data 字段应已移除"


class TestPdfConverterDocling119:
    """P1-6: docling 2.119 页面结构适配 (pages 值为无内容的 PageItem 引用)"""

    def test_convert_document_items_path(self, tmp_path, monkeypatch):
        """2.119 路径: pages 无 cells/items → 从 document 级 texts/tables 按 prov.page_no 归属页面"""
        from src.converters.pdf_converter import PdfConverter

        class FakeProv:
            def __init__(self, page_no):
                self.page_no = page_no

        class FakeTextItem:
            """模拟 docling 2.119 TextItem (label 为枚举, 带 value)"""

            def __init__(self, text, page_no, label=None):
                self.text = text
                self.prov = [FakeProv(page_no)]
                self.label = label

        class FakeTableItem:
            def __init__(self, rows, page_no):
                self._rows = rows
                self.prov = [FakeProv(page_no)]

            def export_to_dataframe(self):
                class FakeDataFrame:
                    def __init__(self, rows):
                        self._rows = rows

                    def iterrows(self):
                        yield from enumerate(self._rows)

                return FakeDataFrame(self._rows)

        class FakePageItem:
            """2.119 PageItem: 仅 size/image/page_no，无 cells/items/tables"""

            pass

        class FakeDoclingDoc:
            pages = {1: FakePageItem(), 2: FakePageItem()}
            texts = [
                FakeTextItem("第一页正文", page_no=1),
                FakeTextItem("第二页标题", page_no=2, label="title"),
                FakeTextItem("第二页正文", page_no=2),
            ]
            tables = [FakeTableItem([["A", "B"], ["1", "2"]], page_no=1)]

        class FakeResult:
            document = FakeDoclingDoc()

        class FakeDoclingConverter:
            def __init__(self, *a, **kw):
                pass

            def convert(self, path):
                return FakeResult()

        _inject_fake_docling(monkeypatch, FakeDoclingConverter)
        fake = tmp_path / "sample.pdf"
        fake.write_bytes(b"%PDF-1.4 fake")

        doc = PdfConverter().convert(str(fake))
        assert doc.metadata.page_count == 2
        assert len(doc.pages) == 2
        # 结构路径: 每页都有 text_frame 元素 (非纯 fallback 空页)
        for page in doc.pages:
            text_frames = [e for e in page.flattened_elements if e.type == "text_frame"]
            assert text_frames, f"page {page.index} 应有 text_frame 元素"
        assert "第一页正文" in doc.pages[0].all_text
        assert "第二页标题" in doc.pages[1].all_text
        # label=title → level 1 / is_title
        page2_titles = [
            p for e in doc.pages[1].flattened_elements for p in e.paragraphs if p.level == 1
        ]
        assert page2_titles and page2_titles[0].text == "第二页标题"
        # 表格归属第 1 页
        tables = [e for e in doc.pages[0].flattened_elements if e.type == "table"]
        assert len(tables) == 1
        assert [c.text for row in tables[0].tables for c in row] == ["A", "B", "1", "2"]

    def test_item_spanning_multiple_pages_assigned_to_all(self, tmp_path, monkeypatch):
        """F3: 多 prov 页归属 — 跨页 item 并入所有涉及页面 (旧实现只取 prov[0])"""
        from src.converters.pdf_converter import PdfConverter

        class FakeProv:
            def __init__(self, page_no):
                self.page_no = page_no

        class FakeTextItem:
            def __init__(self, text, page_nos, label=None):
                self.text = text
                self.prov = [FakeProv(p) for p in page_nos]
                self.label = label

        class FakePageItem:
            pass

        class FakeDoclingDoc:
            pages = {1: FakePageItem(), 2: FakePageItem()}
            texts = [FakeTextItem("跨页文本", page_nos=[1, 2])]
            tables = []

        class FakeResult:
            document = FakeDoclingDoc()

        class FakeDoclingConverter:
            def __init__(self, *a, **kw):
                pass

            def convert(self, path):
                return FakeResult()

        _inject_fake_docling(monkeypatch, FakeDoclingConverter)
        fake = tmp_path / "span.pdf"
        fake.write_bytes(b"%PDF-1.4 fake")

        doc = PdfConverter().convert(str(fake))
        assert "跨页文本" in doc.pages[0].all_text, "跨页 item 应归属第 1 页"
        assert "跨页文本" in doc.pages[1].all_text, "跨页 item 应归属第 2 页"

    def test_convert_orphan_items_go_to_first_page(self, tmp_path, monkeypatch):
        """无 prov 归属的文本并入第一页，不丢内容"""
        from src.converters.pdf_converter import PdfConverter

        class FakeTextItem:
            def __init__(self, text, page_no=None):
                self.text = text
                self.prov = [type("P", (), {"page_no": page_no})()] if page_no else []

        class FakePageItem:
            pass

        class FakeDoclingDoc:
            pages = {1: FakePageItem()}
            texts = [FakeTextItem("孤儿文本", page_no=None)]
            tables = []

        class FakeResult:
            document = FakeDoclingDoc()

        class FakeDoclingConverter:
            def __init__(self, *a, **kw):
                pass

            def convert(self, path):
                return FakeResult()

        _inject_fake_docling(monkeypatch, FakeDoclingConverter)
        fake = tmp_path / "orphan.pdf"
        fake.write_bytes(b"%PDF-1.4 fake")

        doc = PdfConverter().convert(str(fake))
        assert "孤儿文本" in doc.pages[0].all_text

    def test_hf_cache_env_injected_before_instantiation(self, tmp_path, monkeypatch):
        """M5: convert() 在实例化 DoclingConverter 前注入 HF_HUB_CACHE 指向 packages/hf_cache；
        目录存在时额外设 HF_HUB_OFFLINE=1"""
        import os

        import src.converters.pdf_converter as pdf_mod
        from src.converters.pdf_converter import PdfConverter

        hf_cache = tmp_path / "hf_cache"
        monkeypatch.setattr(pdf_mod, "_hf_cache_dir", lambda: hf_cache)
        monkeypatch.delenv("HF_HUB_CACHE", raising=False)
        monkeypatch.delenv("HF_HUB_OFFLINE", raising=False)

        state = {}

        class FakeDoclingDoc:
            pages = None

            def export_to_markdown(self):
                return ""

        class FakeResult:
            document = FakeDoclingDoc()

        class FakeDoclingConverter:
            def __init__(self, *a, **kw):
                state["env_at_init"] = dict(os.environ)

            def convert(self, path):
                return FakeResult()

        _inject_fake_docling(monkeypatch, FakeDoclingConverter)
        fake = tmp_path / "env.pdf"
        fake.write_bytes(b"%PDF-1.4 fake")

        # 状态一: 目录不存在 → 仅 HF_HUB_CACHE，不设 offline
        PdfConverter().convert(str(fake))
        assert state["env_at_init"].get("HF_HUB_CACHE") == str(hf_cache), (
            "HF_HUB_CACHE 应在 DoclingConverter 实例化前注入"
        )
        assert state["env_at_init"].get("HF_HUB_OFFLINE") is None
        assert os.environ.get("HF_HUB_CACHE") == str(hf_cache)

        # 状态二: 目录存在 → HF_HUB_CACHE + HF_HUB_OFFLINE=1
        hf_cache.mkdir(parents=True, exist_ok=True)
        monkeypatch.delenv("HF_HUB_CACHE", raising=False)
        monkeypatch.delenv("HF_HUB_OFFLINE", raising=False)
        state.clear()
        PdfConverter().convert(str(fake))
        assert state["env_at_init"].get("HF_HUB_CACHE") == str(hf_cache)
        assert state["env_at_init"].get("HF_HUB_OFFLINE") == "1", (
            "缓存目录存在时应设 HF_HUB_OFFLINE=1"
        )

    def test_real_conversion_with_docling(self, tmp_path):
        """真实集成测试: docling 可用时真实转换 sample.pdf。

        断言 pages >= 1、文本包含 DocAuditTest、走结构路径 (元素含 text_frame 而非纯 fallback)。
        注意: 评估环境沙箱会拦截 docling-parse 的 C fopen (伪"文件不存在")，真实环境不受影响。
        """
        pytest.importorskip("docling")
        from src.converters.pdf_converter import PdfConverter

        sample_pdf = Path(__file__).parent / "fixtures" / "sample.pdf"
        assert sample_pdf.exists(), f"fixture 缺失: {sample_pdf}"

        doc = PdfConverter().convert(str(sample_pdf))
        assert doc.format == "pdf"
        assert doc.metadata.page_count is not None
        assert doc.metadata.page_count >= 1
        assert len(doc.pages) >= 1
        assert "DocAuditTest" in doc.all_text
        # 结构路径: 有 text_frame 元素，而非最终回退的空页
        text_frames = [
            e for page in doc.pages for e in page.flattened_elements if e.type == "text_frame"
        ]
        assert len(text_frames) >= 1


class TestMarkdownFrontmatterStrict:
    """P1-8: frontmatter 严格判定 (首行独立 --- + 闭合行 + YAML mapping)"""

    def test_proper_frontmatter_still_parsed(self, md_converter, tmp_path):
        """触发: 标准三行式 frontmatter → title/author 提取 + 正文完整"""
        md = "---\ntitle: 测试文档\nauthor: 张三\n---\n\n# 标题\n\n正文内容\n"
        path = tmp_path / "fm.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        assert doc.metadata.title == "测试文档"
        assert doc.metadata.author == "张三"
        assert "正文内容" in doc.all_text

    def test_body_starting_with_dashes_not_swallowed(self, md_converter, tmp_path):
        """不触发: 正文以 '---' 开头 (横向分隔线) → 不吞正文、不崩溃"""
        md = "---\n正文第一段\n---\n正文第二段\n"
        path = tmp_path / "dash.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        assert "正文第一段" in doc.all_text
        assert "正文第二段" in doc.all_text

    def test_dash_separated_sections_not_swallowed(self, md_converter, tmp_path):
        """不触发: '---' 装饰线开头 + 章节分隔 → 各章节完整保留

        (F12 后 '# 章节一' 属注释-only 内容会判为 frontmatter，改用非注释内容验证装饰线不吞正文)
        """
        md = "---\n章节一正文\n\n---\n# 章节二\n"
        path = tmp_path / "sections.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        assert "章节一正文" in doc.all_text
        assert "章节二" in doc.all_text

    def test_non_standalone_dash_line_not_frontmatter(self, md_converter, tmp_path):
        """不触发: 首行 '--- 装饰线' (非独立 ---) → 按正文保留"""
        md = "--- 装饰线\n正文内容\n"
        path = tmp_path / "deco.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        assert "--- 装饰线" in doc.all_text
        assert "正文内容" in doc.all_text

    def test_comment_only_frontmatter_detected(self, md_converter, tmp_path):
        """F12: --- 开头且闭合、内容仅注释/空白 → 判定为 frontmatter (不当作正文)"""
        md = "---\n# 生成日期: 2026-08-19\n# 来源: 内部\n---\n\n正文内容\n"
        path = tmp_path / "cmt.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        assert doc.metadata.title == "cmt"  # frontmatter 无 title → 回退文件名
        assert "生成日期" not in doc.all_text, "注释-only frontmatter 应被消费，不进入正文"
        assert "正文内容" in doc.all_text


class TestMarkdownTableAndHeadingFixes:
    """P1-8: 表格分隔行误判 + 标题前缀空格一致"""

    def test_separator_row_still_skipped(self, md_converter, tmp_path):
        """触发: 标准分隔行 (|---|) 仍被跳过"""
        md = "| 名称 | 数值 |\n| --- | --- |\n| A | 1 |\n"
        path = tmp_path / "sep.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        tables = [e for p in doc.pages for e in p.flattened_elements if e.type == "table"]
        assert len(tables) >= 1
        cells = [c.text for row in tables[0].tables for c in row]
        assert "---" not in cells
        assert "A" in cells

    def test_dash_data_row_not_treated_as_separator(self, md_converter, tmp_path):
        """不触发: 内容为纯 '-' 的数据行不得被当分隔行跳过"""
        md = "| 名称 | 数值 |\n| --- | --- |\n| - | - |\n| A | 1 |\n"
        path = tmp_path / "dashrow.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        tables = [e for p in doc.pages for e in p.flattened_elements if e.type == "table"]
        assert len(tables) >= 1
        cells = [c.text for row in tables[0].tables for c in row]
        assert "-" in cells, f"'-' 数据行被误当分隔行跳过: {cells}"
        assert "A" in cells and "1" in cells

    def test_inline_regexes_precompiled_module_level(self):
        """F8: 列表/分隔行/分隔格正则预编译为模块级常量 (与 HEADING_RE 同风格)"""
        import re as _re

        from src.converters import md_converter as mdc

        for name in ("_LIST_MARKER_RE", "_HR_RE", "_SEPARATOR_CELL_RE"):
            const = getattr(mdc, name, None)
            assert const is not None, f"{name} 应存在"
            assert isinstance(const, _re.Pattern), f"{name} 应为编译后 Pattern"

    def test_table_rows_continuous_after_separator(self, md_converter, tmp_path):
        """F2: 分隔行跳过不影响后续行的 row 索引连续 (FMT-008 位置对齐数据源)"""
        md = "| 名称 | 数值 |\n| --- | --- |\n| A | 1 |\n| B | 2 |\n"
        path = tmp_path / "cont.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        tables = [e for p in doc.pages for e in p.flattened_elements if e.type == "table"]
        assert len(tables) >= 1
        rows = tables[0].tables
        # 表头 + 2 数据行 = 3 行；分隔行跳过但 row 索引连续无空洞 (旧实现为 [2,2,3,3])
        assert [c.row for row in rows for c in row] == [0, 0, 1, 1, 2, 2], (
            f"分隔行后 row 索引应连续, got {[c.row for row in rows for c in row]}"
        )

    def test_heading_with_leading_spaces_detected(self, md_converter, tmp_path):
        """触发: 0-3 空格前缀的标题被识别为标题 (与页面分割正则 s{0,3} 一致)"""
        md = "# 一级\n\n正文段落\n   #### 缩进四级标题\n"
        path = tmp_path / "indent.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        levels = [p.level for p in doc.all_paragraphs if p.level is not None]
        assert 4 in levels, f"带 3 空格前缀的 #### 标题应识别为 level 4: {levels}"
        assert "缩进四级标题" in doc.all_text

    def test_four_space_indent_not_heading(self, md_converter, tmp_path):
        """不触发: 4 空格缩进 (代码) 不误判为标题"""
        md = "# 标题\n\n    #### 代码缩进\n"
        path = tmp_path / "codeindent.md"
        path.write_text(md, encoding="utf-8")
        doc = md_converter.convert(str(path))
        levels = [p.level for p in doc.all_paragraphs if p.level is not None]
        assert 4 not in levels, f"4 空格缩进不应识别为标题: {levels}"
        assert "#### 代码缩进" in doc.all_text
